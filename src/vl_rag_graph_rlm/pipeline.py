"""Unified Multimodal RAG Pipeline - Single interface for VL-RAG-Graph-RLM.

This module provides a high-level, unified API that combines:
- Paddle OCR for PDF/image extraction
- Qwen3-VL for multimodal embeddings and reranking
- Hybrid search (dense + keyword with RRF)
- VLRAGGraphRLM for recursive reasoning
- Cheap SOTA LLMs via OpenRouter

Example:
    >>> from vl_rag_graph_rlm import MultimodalRAGPipeline
    
    # Initialize unified pipeline
    pipeline = MultimodalRAGPipeline(
        llm_provider="openrouter",
        llm_model="kimi/kimi-k2.5",
        embedding_model="Qwen/Qwen3-VL-Embedding-2B",
        use_reranker=True
    )
    
    # Process PDF with images
    pipeline.add_pdf("research_paper.pdf", extract_images=True)
    
    # Query with automatic retrieval and reasoning
    result = pipeline.query("Explain Figure 3 and its implications")
    print(result.answer)
    print(f"Cost: ${result.total_cost:.4f}")
"""

import os
import logging
from typing import List, Dict, Any, Optional, Union, Callable
from dataclasses import dataclass, field
from pathlib import Path

from vl_rag_graph_rlm.core import VLRAGGraphRLM
from vl_rag_graph_rlm.rag.multimodal_store import MultimodalVectorStore

try:
    from vl_rag_graph_rlm.rag.qwen3vl import (
        Qwen3VLEmbeddingProvider, 
        Qwen3VLRerankerProvider,
        create_qwen3vl_embedder,
        create_qwen3vl_reranker
    )
    HAS_QWEN3VL = True
except ImportError:
    HAS_QWEN3VL = False
    Qwen3VLEmbeddingProvider = None  # type: ignore
    Qwen3VLRerankerProvider = None  # type: ignore
    create_qwen3vl_embedder = None  # type: ignore
    create_qwen3vl_reranker = None  # type: ignore

from vl_rag_graph_rlm.rag import SearchResult, ReciprocalRankFusion, MultiFactorReranker
from vl_rag_graph_rlm.types import VLRAGGraphRLMChatCompletion as RLMChatCompletion

logger = logging.getLogger("vl_rag_graph_rlm.pipeline")


@dataclass
class PipelineResult:
    """Result from a pipeline query."""
    answer: str
    sources: List[Dict[str, Any]] = field(default_factory=list)
    retrieved_content: List[str] = field(default_factory=list)
    images: List[str] = field(default_factory=list)
    total_cost: float = 0.0
    execution_time: float = 0.0
    llm_calls: int = 0
    iterations: int = 0
    depth: int = 0


class MultimodalRAGPipeline:
    """Unified pipeline for multimodal document analysis.
    
    Combines OCR, multimodal embeddings, hybrid search, reranking, 
    and recursive LLM reasoning into a single interface.
    
    Example:
        >>> pipeline = MultimodalRAGPipeline(
        ...     llm_provider="openrouter",
        ...     llm_model="kimi/kimi-k2.5",
        ...     recursive_model="google/gemini-2.0-flash-001"
        ... )
        ...
        >>> # Add documents
        >>> pipeline.add_pdf("doc.pdf", extract_images=True)
        >>> pipeline.add_image("diagram.png")
        ...
        >>> # Query
        >>> result = pipeline.query("What does the diagram show?")
        >>> print(result.answer)
    """
    
    def __init__(
        self,
        # LLM Configuration
        llm_provider: str = "openrouter",
        llm_model: Optional[str] = None,
        recursive_model: Optional[str] = None,
        llm_api_key: Optional[str] = None,
        # Embedding Configuration
        embedding_model: str = "Qwen/Qwen3-VL-Embedding-2B",
        embedding_device: Optional[str] = None,
        # Reranker Configuration
        use_reranker: bool = True,
        reranker_model: str = "Qwen/Qwen3-VL-Reranker-2B",
        # Search Configuration
        top_k: int = 10,
        rerank_top_k: int = 100,
        use_hybrid_search: bool = True,
        keyword_weight: float = 1.0,
        dense_weight: float = 4.0,
        # RLM Configuration
        max_depth: int = 3,
        max_iterations: int = 10,
        temperature: float = 0.0,
        # Storage
        storage_path: Optional[str] = None,
        # OCR Configuration
        ocr_language: str = "en",
        extract_tables: bool = True,
        extract_images: bool = True,
        # Misc
        verbose: bool = False
    ):
        """Initialize the unified multimodal RAG pipeline.
        
        Args:
            llm_provider: LLM provider (openrouter, openai, anthropic, etc.)
            llm_model: LLM model name (defaults to provider's cheap SOTA)
            recursive_model: Cheaper model for recursive calls
            llm_api_key: API key for LLM provider
            embedding_model: Qwen3-VL embedding model name
            embedding_device: Device for embeddings (cuda/cpu/mps)
            use_reranker: Whether to use Qwen3-VL reranker
            reranker_model: Qwen3-VL reranker model name
            top_k: Number of final results to retrieve
            rerank_top_k: Number of results to rerank
            use_hybrid_search: Whether to combine dense + keyword search
            keyword_weight: Weight for keyword search in hybrid
            dense_weight: Weight for dense search in hybrid
            max_depth: Max recursion depth for RLM
            max_iterations: Max iterations per RLM call
            temperature: Sampling temperature
            storage_path: Path to persist document store
            ocr_language: Language for OCR (en, ch, etc.)
            extract_tables: Whether to extract tables from PDFs
            extract_images: Whether to extract images from PDFs
            verbose: Whether to log detailed progress
        """
        self.verbose = verbose
        if verbose:
            logging.basicConfig(level=logging.INFO)
        
        # Initialize embedding provider
        logger.info(f"Loading embedding model: {embedding_model}")
        self.embedding_provider = create_qwen3vl_embedder(
            model_name=embedding_model,
            device=embedding_device
        )
        
        # Initialize reranker if requested
        self.reranker = None
        if use_reranker:
            logger.info(f"Loading reranker model: {reranker_model}")
            self.reranker = create_qwen3vl_reranker(
                model_name=reranker_model,
                device=embedding_device
            )
        
        # Initialize vector store
        self.store = MultimodalVectorStore(
            embedding_provider=self.embedding_provider,
            storage_path=storage_path,
            use_qwen_reranker=use_reranker,
            reranker_provider=self.reranker
        )
        
        # Initialize RLM
        logger.info(f"Initializing RLM with {llm_provider}")
        self.rlm = VLRAGGraphRLM(
            provider=llm_provider,
            model=llm_model,
            recursive_model=recursive_model,
            api_key=llm_api_key,
            max_depth=max_depth,
            max_iterations=max_iterations,
            temperature=temperature
        )
        
        # Search configuration
        self.top_k = top_k
        self.rerank_top_k = rerank_top_k
        self.use_hybrid_search = use_hybrid_search
        self.keyword_weight = keyword_weight
        self.dense_weight = dense_weight
        
        # OCR configuration
        self.ocr_language = ocr_language
        self.extract_tables = extract_tables
        self.extract_images = extract_images
        
        # Hybrid search components
        self.rrf = ReciprocalRankFusion(k=60)
        self.keyword_reranker = MultiFactorReranker()
        
        logger.info("Pipeline initialized successfully")
    
    def add_pdf(
        self,
        pdf_path: Union[str, Path],
        pages: Optional[List[int]] = None,
        metadata: Optional[Dict[str, Any]] = None,
        extract_images: Optional[bool] = None,
        extract_tables: Optional[bool] = None,
        use_ocr: bool = True
    ) -> List[str]:
        """Add a PDF document to the pipeline.
        
        Extracts text, images, and tables using Paddle OCR or pdf2image.
        
        Args:
            pdf_path: Path to PDF file
            pages: Specific pages to extract (None = all pages)
            metadata: Additional metadata for the document
            extract_images: Whether to extract images (default: from init)
            extract_tables: Whether to extract tables (default: from init)
            use_ocr: Whether to use Paddle OCR for layout analysis
            
        Returns:
            List of document IDs added
        """
        pdf_path = Path(pdf_path)
        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF not found: {pdf_path}")
        
        extract_images = extract_images if extract_images is not None else self.extract_images
        extract_tables = extract_tables if extract_tables is not None else self.extract_tables
        metadata = metadata or {}
        
        doc_ids = []
        
        if use_ocr:
            # Use Paddle OCR for advanced layout extraction
            doc_ids = self._add_pdf_with_paddle_ocr(
                pdf_path, pages, metadata, extract_images, extract_tables
            )
        else:
            # Use simple pdf2image approach
            doc_ids = self.store.add_pdf(
                str(pdf_path),
                pages=pages,
                metadata=metadata
            )
        
        logger.info(f"Added {len(doc_ids)} documents from {pdf_path}")
        return doc_ids
    
    def add_presentation(
        self,
        pptx_path: Union[str, Path],
        slides: Optional[List[int]] = None,
        metadata: Optional[Dict[str, Any]] = None,
        extract_images: Optional[bool] = None,
        extract_notes: bool = True,
        extract_tables: bool = True
    ) -> List[str]:
        """Add a PowerPoint presentation to the pipeline.
        
        Extracts text, images, tables, and speaker notes from .pptx files.
        Each slide is processed as a separate document with slide images
        extracted for multimodal retrieval.
        
        Requirements:
            pip install python-pptx Pillow
        
        Args:
            pptx_path: Path to PowerPoint file (.pptx)
            slides: Specific slide numbers to extract (None = all slides)
            metadata: Additional metadata for the document
            extract_images: Whether to extract slide images (default: from init)
            extract_notes: Whether to extract speaker notes
            extract_tables: Whether to extract tables from slides
            
        Returns:
            List of document IDs added
        """
        try:
            from pptx import Presentation
            from PIL import Image
        except ImportError:
            raise ImportError(
                "PowerPoint support requires python-pptx and Pillow. "
                "Install with: pip install python-pptx Pillow"
            )
        
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
        if pptx_path.suffix.lower() != '.pptx':
            raise ValueError(f"File must be a .pptx file, got: {pptx_path.suffix}")
        
        extract_images = extract_images if extract_images is not None else self.extract_images
        metadata = metadata or {}
        doc_ids = []
        
        logger.info(f"Processing PowerPoint: {pptx_path}")
        prs = Presentation(str(pptx_path))
        
        # Determine which slides to process
        slides_to_process = slides or range(1, len(prs.slides) + 1)
        
        for slide_num in slides_to_process:
            if not (1 <= slide_num <= len(prs.slides)):
                logger.warning(f"Slide {slide_num} out of range, skipping")
                continue
            
            slide = prs.slides[slide_num - 1]  # 0-indexed
            slide_content_parts = []
            slide_metadata = metadata.copy()
            slide_metadata.update({
                "source": str(pptx_path),
                "slide_number": slide_num,
                "total_slides": len(prs.slides),
                "type": "presentation_slide"
            })
            
            # Extract slide layout name
            if slide.slide_layout.name:
                slide_metadata["layout"] = slide.slide_layout.name
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content_parts.append(shape.text.strip())
                
                # Extract tables
                if extract_tables and shape.has_table:
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(" | ".join(row_data))
                    if table_data:
                        slide_content_parts.append("\n[Table]\n" + "\n".join(table_data))
                        slide_metadata["has_table"] = True
            
            # Extract speaker notes
            if extract_notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                if notes_text_frame.text.strip():
                    slide_content_parts.append(f"\n[Speaker Notes]\n{notes_text_frame.text.strip()}")
                    slide_metadata["has_notes"] = True
            
            # Combine slide content
            slide_content = "\n\n".join(slide_content_parts)
            
            # Add slide text as document
            if slide_content.strip():
                doc_id = self.store.add_text(
                    content=slide_content,
                    metadata=slide_metadata
                )
                doc_ids.append(doc_id)
                logger.debug(f"Added slide {slide_num} text content")
            
            # Extract and process slide as image if enabled
            if extract_images:
                try:
                    # Save slide as image using slide rendering
                    # Note: python-pptx doesn't directly support slide export
                    # We'll extract embedded images instead
                    for shape_idx, shape in enumerate(slide.shapes):
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            try:
                                image = shape.image
                                image_bytes = image.blob
                                
                                # Save temporarily
                                temp_img_path = f"/tmp/{pptx_path.stem}_slide{slide_num}_img{shape_idx}.png"
                                with open(temp_img_path, "wb") as f:
                                    f.write(image_bytes)
                                
                                # Add to store with description
                                img_metadata = metadata.copy()
                                img_metadata.update({
                                    "source": str(pptx_path),
                                    "slide_number": slide_num,
                                    "type": "presentation_image",
                                    "shape_index": shape_idx
                                })
                                
                                doc_id = self.store.add_image(
                                    image_path=temp_img_path,
                                    description=f"Image from slide {slide_num}",
                                    metadata=img_metadata
                                )
                                doc_ids.append(doc_id)
                            except Exception as e:
                                logger.warning(f"Failed to extract image from slide {slide_num}: {e}")
                except Exception as e:
                    logger.warning(f"Failed to process images for slide {slide_num}: {e}")
        
        logger.info(f"Added {len(doc_ids)} documents from PowerPoint: {pptx_path}")
        return doc_ids
    
    def _add_pdf_with_paddle_ocr(
        self,
        pdf_path: Path,
        pages: Optional[List[int]],
        metadata: Dict[str, Any],
        extract_images: bool,
        extract_tables: bool
    ) -> List[str]:
        """Add PDF using Paddle OCR for layout analysis."""
        try:
            from paddleocr import PaddleOCR
            import fitz  # PyMuPDF
        except ImportError:
            logger.warning("PaddleOCR or PyMuPDF not installed. Using simple PDF extraction.")
            return self.store.add_pdf(str(pdf_path), pages=pages, metadata=metadata)
        
        doc_ids = []
        
        # Initialize Paddle OCR
        ocr = PaddleOCR(
            use_angle_cls=True,
            lang=self.ocr_language,
            use_gpu=False  # Set to True if GPU available
        )
        
        # Open PDF
        doc = fitz.open(str(pdf_path))
        pages_to_process = pages or range(1, len(doc) + 1)
        
        for page_num in pages_to_process:
            if isinstance(page_num, int) and page_num <= len(doc):
                page = doc[page_num - 1]  # 0-indexed in PyMuPDF
                
                # Extract text with OCR
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                from PIL import Image
                import io
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                
                # Run OCR
                ocr_result = ocr.ocr(pix.tobytes("png"), cls=True)
                
                # Extract text
                texts = []
                if ocr_result[0]:
                    for line in ocr_result[0]:
                        if line:
                            texts.append(line[1][0])  # Text content
                
                content = "\n".join(texts)
                
                # Create metadata
                page_metadata = metadata.copy()
                page_metadata.update({
                    "source": str(pdf_path),
                    "page": page_num,
                    "total_pages": len(doc),
                    "type": "pdf_page"
                })
                
                # Add as multimodal document
                if content.strip():
                    doc_id = self.store.add_text(
                        content=content,
                        metadata=page_metadata
                    )
                    doc_ids.append(doc_id)
                
                # Extract images if requested
                if extract_images:
                    image_list = page.get_images()
                    for img_idx, img in enumerate(image_list):
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        
                        # Save temporarily
                        temp_img_path = f"/tmp/{pdf_path.stem}_p{page_num}_img{img_idx}.png"
                        with open(temp_img_path, "wb") as f:
                            f.write(image_bytes)
                        
                        # Add to store
                        img_metadata = metadata.copy()
                        img_metadata.update({
                            "source": str(pdf_path),
                            "page": page_num,
                            "type": "extracted_image",
                            "image_index": img_idx
                        })
                        
                        doc_id = self.store.add_image(
                            image_path=temp_img_path,
                            description=f"Image {img_idx} from page {page_num}",
                            metadata=img_metadata
                        )
                        doc_ids.append(doc_id)
        
        doc.close()
        return doc_ids
    
    def add_image(
        self,
        image_path: Union[str, Path],
        description: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None
    ) -> str:
        """Add an image document to the pipeline.
        
        Args:
            image_path: Path to image file
            description: Optional text description
            metadata: Additional metadata
            
        Returns:
            Document ID
        """
        metadata = metadata or {}
        return self.store.add_image(
            image_path=str(image_path),
            description=description,
            metadata=metadata
        )
    
    def add_text(
        self,
        content: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> str:
        """Add a text document to the pipeline.
        
        Args:
            content: Text content
            metadata: Additional metadata
            
        Returns:
            Document ID
        """
        metadata = metadata or {}
        return self.store.add_text(content=content, metadata=metadata)
    
    def query(
        self,
        query: str,
        top_k: Optional[int] = None,
        use_reranking: Optional[bool] = None,
        return_sources: bool = True,
        system_prompt: Optional[str] = None
    ) -> PipelineResult:
        """Query the pipeline with automatic retrieval and reasoning.
        
        This is the main interface - it retrieves relevant documents,
        optionally reranks them, and uses the RLM to generate an answer.
        
        Args:
            query: User query (can include references to images/figures)
            top_k: Number of documents to retrieve (default: from init)
            use_reranking: Whether to rerank results (default: from init)
            return_sources: Whether to include source documents in result
            system_prompt: Optional custom system prompt for RLM
            
        Returns:
            PipelineResult with answer and metadata
        """
        import time
        start_time = time.time()
        
        top_k = top_k or self.top_k
        use_reranking = use_reranking if use_reranking is not None else (self.reranker is not None)
        
        logger.info(f"Querying: {query[:50]}...")
        
        # Step 1: Retrieve relevant documents
        if self.use_hybrid_search:
            results = self._hybrid_search(query, top_k=top_k * 2 if use_reranking else top_k)
        else:
            results = self.store.search(query, top_k=top_k * 2 if use_reranking else top_k)
        
        logger.info(f"Retrieved {len(results)} documents")
        
        # Step 2: Rerank if enabled
        if use_reranking and self.reranker:
            results = self._rerank_results(query, results, top_k=top_k)
            logger.info(f"Reranked to {len(results)} documents")
        else:
            results = results[:top_k]
        
        # Step 3: Build context from retrieved documents
        context_parts = []
        sources = []
        images = []
        
        for i, result in enumerate(results):
            doc = self.store.get(result.id)
            if not doc:
                continue
            
            # Add to context
            context_text = f"[Document {i+1}]\n{doc.content}\n"
            if doc.metadata:
                context_text += f"Metadata: {doc.metadata}\n"
            context_parts.append(context_text)
            
            # Track sources
            if return_sources:
                sources.append({
                    "id": doc.id,
                    "content": doc.content[:200] + "..." if len(doc.content) > 200 else doc.content,
                    "metadata": doc.metadata,
                    "score": result.composite_score if hasattr(result, 'composite_score') else result.semantic_score
                })
            
            # Track images
            if doc.image_path:
                images.append(doc.image_path)
        
        context = "\n\n".join(context_parts)
        
        # Step 4: Generate answer with RLM
        logger.info("Generating answer with RLM...")
        
        # Enhance query with context about available images
        enhanced_query = query
        if images:
            enhanced_query += f"\n\nAvailable images: {len(images)} image(s) referenced in the retrieved documents."
        
        rlm_result = self.rlm.completion(
            query=enhanced_query,
            context=context
        )
        
        execution_time = time.time() - start_time
        
        # Get usage stats
        usage = rlm_result.usage_summary
        total_cost = usage.total_cost if usage else 0.0
        
        return PipelineResult(
            answer=rlm_result.response,
            sources=sources if return_sources else [],
            retrieved_content=[doc.content for doc in [self.store.get(r.id) for r in results] if doc],
            images=images,
            total_cost=total_cost,
            execution_time=execution_time,
            llm_calls=self.rlm.stats.get('llm_calls', 0),
            iterations=self.rlm.stats.get('iterations', 0),
            depth=self.rlm.stats.get('depth', 0)
        )
    
    def _hybrid_search(self, query: str, top_k: int = 10) -> List[SearchResult]:
        """Perform hybrid search combining dense and keyword search."""
        # Dense search
        dense_results = self.store.search(query, top_k=top_k)
        
        # Keyword search (simple implementation)
        keyword_results = self._keyword_search(query, top_k=top_k)
        
        # Fuse results with RRF
        if keyword_results:
            fused = self.rrf.fuse(
                [dense_results, keyword_results],
                weights=[self.dense_weight, self.keyword_weight]
            )
            return fused
        
        return dense_results
    
    def _keyword_search(self, query: str, top_k: int = 10) -> List[SearchResult]:
        """Simple keyword-based search."""
        import re
        
        query_terms = set(re.findall(r'\w+', query.lower()))
        results = []
        
        for doc in self.store.documents.values():
            content_lower = doc.content.lower()
            matches = sum(1 for term in query_terms if term in content_lower)
            
            if matches > 0:
                score = matches / len(query_terms) if query_terms else 0
                results.append(SearchResult(
                    id=doc.id,
                    content=doc.content,
                    metadata=doc.metadata,
                    keyword_score=score,
                    composite_score=score
                ))
        
        # Sort by keyword score
        results.sort(key=lambda x: x.keyword_score, reverse=True)
        return results[:top_k]
    
    def _rerank_results(
        self,
        query: str,
        results: List[SearchResult],
        top_k: int = 10
    ) -> List[SearchResult]:
        """Rerank results using Qwen3-VL reranker."""
        if not self.reranker:
            return results[:top_k]
        
        # Format documents for reranker
        documents = []
        for result in results:
            doc = self.store.get(result.id)
            if doc:
                doc_dict = {"text": doc.content}
                if doc.image_path:
                    doc_dict["image"] = doc.image_path
                documents.append(doc_dict)
        
        if not documents:
            return results[:top_k]
        
        # Rerank
        reranked_indices = self.reranker.rerank(
            query={"text": query},
            documents=documents
        )
        
        # Reorder results based on reranking
        reordered = []
        for idx, score in reranked_indices[:top_k]:
            result = results[idx]
            result.composite_score = score * 100
            reordered.append(result)
        
        return reordered
    
    def search(
        self,
        query: str,
        query_image: Optional[Union[str, Path]] = None,
        top_k: int = 10,
        use_reranking: bool = True
    ) -> List[SearchResult]:
        """Search documents without generating an answer.
        
        Args:
            query: Text query
            query_image: Optional image query
            top_k: Number of results
            use_reranking: Whether to rerank
            
        Returns:
            List of search results
        """
        # Multimodal search if image provided
        if query_image:
            results = self.store.search_multimodal(
                text=query,
                image=str(query_image),
                top_k=top_k * 2 if use_reranking else top_k
            )
        else:
            results = self.store.search(
                query,
                top_k=top_k * 2 if use_reranking else top_k
            )
        
        # Rerank if enabled
        if use_reranking and self.reranker:
            results = self._rerank_results(query, results, top_k=top_k)
        else:
            results = results[:top_k]
        
        return results
    
    def get_document(self, doc_id: str) -> Optional[Dict[str, Any]]:
        """Get a document by ID.
        
        Args:
            doc_id: Document ID
            
        Returns:
            Document data or None
        """
        doc = self.store.get(doc_id)
        if doc:
            return {
                "id": doc.id,
                "content": doc.content,
                "metadata": doc.metadata,
                "image_path": doc.image_path,
                "video_path": doc.video_path
            }
        return None
    
    def delete_document(self, doc_id: str) -> bool:
        """Delete a document by ID.
        
        Args:
            doc_id: Document ID
            
        Returns:
            True if deleted, False if not found
        """
        return self.store.delete(doc_id)
    
    def save(self):
        """Persist the document store to disk."""
        self.store._save()
        logger.info(f"Saved document store to {self.store.storage_path}")
    
    def load(self):
        """Load the document store from disk."""
        self.store._load()
        logger.info(f"Loaded document store from {self.store.storage_path}")
    
    @property
    def stats(self) -> Dict[str, Any]:
        """Get pipeline statistics."""
        return {
            "documents": len(self.store.documents),
            "storage_path": self.store.storage_path,
            "embedding_model": getattr(self.embedding_provider, 'model_name_or_path', 'unknown'),
            "reranker_enabled": self.reranker is not None,
            "rlm_stats": self.rlm.stats
        }


# Convenience function for quick usage
def create_pipeline(
    llm_provider: str = "openrouter",
    llm_model: Optional[str] = None,
    embedding_model: str = "Qwen/Qwen3-VL-Embedding-2B",
    use_reranker: bool = True,
    storage_path: Optional[str] = None,
    **kwargs
) -> MultimodalRAGPipeline:
    """Create a preconfigured multimodal RAG pipeline.
    
    Args:
        llm_provider: LLM provider
        llm_model: LLM model (defaults to cheap SOTA for provider)
        embedding_model: Embedding model name
        use_reranker: Whether to use reranker
        storage_path: Path for persistence
        **kwargs: Additional configuration
        
    Returns:
        Configured MultimodalRAGPipeline
    """
    return MultimodalRAGPipeline(
        llm_provider=llm_provider,
        llm_model=llm_model,
        embedding_model=embedding_model,
        use_reranker=use_reranker,
        storage_path=storage_path,
        **kwargs
    )
