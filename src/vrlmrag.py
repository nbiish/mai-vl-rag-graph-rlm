#!/usr/bin/env python3
"""
vrlmrag — VL-RAG-Graph-RLM CLI

Full multimodal document analysis pipeline:
  1. VL:       Qwen3-VL multimodal embeddings (text + images)
  2. RAG:      Hybrid search (dense + keyword) with RRF fusion
  3. Reranker:  Qwen3-VL cross-attention reranking
  4. Graph:    Knowledge graph extraction via RLM
  5. RLM:      Recursive Language Model with REPL sandbox
  6. Report:   Markdown report generation

Usage:
    vrlmrag document.pptx                          # auto: uses provider hierarchy
    vrlmrag --provider sambanova document.pptx     # explicit provider
    vrlmrag --provider nebius doc.pdf -o report.md  # with output file
    vrlmrag ./folder -q "Summarize key concepts"   # auto + custom query
    vrlmrag --show-hierarchy                       # see fallback order
    vrlmrag --list-providers                       # see all providers

Collections (named persistent knowledge stores):
    vrlmrag -c research --add ./papers/            # add docs to collection
    vrlmrag -c research -q "Key findings?"         # query a collection
    vrlmrag -c research -c code -q "How?"          # blend multiple collections
    vrlmrag -c research -i                         # interactive w/ collection
    vrlmrag --collection-list                      # list all collections
    vrlmrag -c research --collection-info          # show collection details

Provider Hierarchy:
    When --provider is omitted (or set to 'auto'), providers are tried in order:
    sambanova → nebius → groq → cerebras → zai → zenmux → openrouter → ...
    Only providers with valid API keys are attempted. Configurable via
    PROVIDER_HIERARCHY env var in .env.
"""

__version__ = "0.1.1"

import json
import os
import sys
import argparse
import time
from pathlib import Path
from datetime import datetime
from typing import List, Optional, Dict, Any

# Load environment variables from .env file in project root
from dotenv import load_dotenv

# Find project root (where .env is located)
project_root = Path(__file__).parent.parent
env_file = project_root / ".env"
if env_file.exists():
    load_dotenv(dotenv_path=env_file)
else:
    # Try loading from current directory as fallback
    load_dotenv()

# Add src to path
sys.path.insert(0, str(Path(__file__).parent))

from vl_rag_graph_rlm import VLRAGGraphRLM
from vl_rag_graph_rlm.clients.hierarchy import (
    get_hierarchy,
    get_available_providers,
    resolve_auto_provider,
    PROVIDER_KEY_MAP,
)
from vl_rag_graph_rlm.rag import (
    CompositeReranker,
    SearchResult,
    ReciprocalRankFusion,
    MultiFactorReranker,
)

# Qwen3-VL imports (requires torch, transformers>=5.1.0)
try:
    from vl_rag_graph_rlm.rag.qwen3vl import (
        Qwen3VLEmbeddingProvider,
        create_qwen3vl_embedder,
        MultimodalDocument,
    )
    from vl_rag_graph_rlm.rag.multimodal_store import MultimodalVectorStore

    HAS_QWEN3VL = True
except ImportError:
    HAS_QWEN3VL = False

# FlashRank lightweight reranker (~34 MB, ONNX cross-encoder)
try:
    from vl_rag_graph_rlm.rag.flashrank_reranker import (
        create_flashrank_reranker,
    )
    HAS_FLASHRANK = True
except ImportError:
    HAS_FLASHRANK = False

# API-based embedding provider (OpenRouter + ZenMux omni)
try:
    from vl_rag_graph_rlm.rag.api_embedding import create_api_embedder
    HAS_API_EMBEDDING = True
except ImportError:
    HAS_API_EMBEDDING = False

# Text-only embedding provider (Qwen3-Embedding, ~1.2 GB)
try:
    from vl_rag_graph_rlm.rag.text_embedding import create_text_embedder
    HAS_TEXT_EMBEDDING = True
except ImportError:
    HAS_TEXT_EMBEDDING = False

# Parakeet ASR transcription (local audio → text, ~0.6 GB)
try:
    from vl_rag_graph_rlm.rag.parakeet import create_parakeet_transcriber
    HAS_PARAKEET = True
except ImportError:
    HAS_PARAKEET = False

from vl_rag_graph_rlm.collections import (
    collection_exists,
    create_collection,
    load_collection_meta,
    save_collection_meta,
    list_collections as _list_collections_meta,
    delete_collection,
    record_source,
    load_kg as collection_load_kg,
    save_kg as collection_save_kg,
    merge_kg as collection_merge_kg,
    _embeddings_path as collection_embeddings_path,
    _kg_path as collection_kg_path,
    _collection_dir,
    COLLECTIONS_ROOT,
)
from vl_rag_graph_rlm.local_model_lock import (
    local_model_lock,
    lock_status,
    is_local_provider,
)

# Optional PPTX support
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


# Audio/video extensions handled by the media pipeline
_AUDIO_EXTENSIONS = {".wav", ".mp3", ".flac", ".m4a", ".ogg", ".opus"}
_VIDEO_EXTENSIONS = {".mp4", ".mkv", ".webm", ".avi", ".mov"}
_MEDIA_EXTENSIONS = _AUDIO_EXTENSIONS | _VIDEO_EXTENSIONS

# All supported file extensions for manifest scanning
_ALL_SUPPORTED_EXTENSIONS = (
    {".txt", ".md", ".pdf", ".pptx", ".docx"}
    | _MEDIA_EXTENSIONS
    | {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg"}
)


def _load_manifest(store_dir: Path) -> dict:
    """Load the file manifest (tracks indexed files + mtimes)."""
    manifest_file = store_dir / "manifest.json"
    if manifest_file.exists():
        try:
            return json.loads(manifest_file.read_text(encoding="utf-8"))
        except Exception:
            return {}
    return {}


def _save_manifest(store_dir: Path, manifest: dict) -> None:
    """Persist the file manifest."""
    store_dir.mkdir(parents=True, exist_ok=True)
    (store_dir / "manifest.json").write_text(
        json.dumps(manifest, indent=2), encoding="utf-8",
    )


def _scan_supported_files(input_path: Path) -> dict[str, float]:
    """Scan input_path for supported files, returning {path: mtime}."""
    files: dict[str, float] = {}
    if input_path.is_file():
        if input_path.suffix.lower() in _ALL_SUPPORTED_EXTENSIONS:
            files[str(input_path)] = input_path.stat().st_mtime
    elif input_path.is_dir():
        for f in sorted(input_path.rglob("*")):
            if f.is_file() and f.suffix.lower() in _ALL_SUPPORTED_EXTENSIONS:
                parts = f.relative_to(input_path).parts
                if any(p.startswith(".") for p in parts):
                    continue
                files[str(f)] = f.stat().st_mtime
    return files


def _detect_file_changes(
    current_files: dict[str, float], manifest: dict,
) -> tuple[list[str], list[str]]:
    """Compare current files against manifest.

    Returns:
        (new_or_modified, deleted) — lists of file paths
    """
    indexed = manifest.get("files", {})
    new_or_modified = [
        fpath for fpath, mtime in current_files.items()
        if indexed.get(fpath) is None or mtime > indexed[fpath]
    ]
    deleted = [f for f in indexed if f not in current_files]
    return new_or_modified, deleted


class DocumentProcessor:
    """Process various document types for VL-RAG-Graph-RLM."""

    def __init__(self, transcription_provider=None, use_api: bool = False):
        self.supported_extensions = {".txt", ".md", ".pdf", ".pptx", ".docx"} | _MEDIA_EXTENSIONS
        self.transcription_provider = transcription_provider
        self.use_api = use_api

    def process_path(self, path: str) -> List[dict]:
        """Process a file or folder path."""
        path_obj = Path(path)

        if path_obj.is_file():
            return [self.process_file(path_obj)]
        elif path_obj.is_dir():
            return self.process_folder(path_obj)
        else:
            raise ValueError(f"Path not found: {path}")

    def process_file(self, file_path: Path) -> dict:
        """Process a single file."""
        ext = file_path.suffix.lower()

        if ext == ".pptx" and HAS_PPTX:
            return self._process_pptx(file_path)
        elif ext == ".pdf":
            return self._process_pdf(file_path)
        elif ext == ".docx":
            return self._process_docx(file_path)
        elif ext == ".csv":
            return self._process_csv(file_path)
        elif ext in {".xlsx", ".xls"}:
            return self._process_excel(file_path)
        elif ext in {".txt", ".md"}:
            return self._process_text(file_path)
        elif ext in _MEDIA_EXTENSIONS:
            return self._process_media(file_path)
        else:
            return {
                "type": "unsupported",
                "path": str(file_path),
                "content": f"File type {ext} not yet supported",
                "chunks": [],
            }

    def process_folder(self, folder_path: Path) -> List[dict]:
        """Process all supported files in a folder."""
        results = []
        for file_path in folder_path.rglob("*"):
            if (
                file_path.is_file()
                and file_path.suffix.lower() in self.supported_extensions
            ):
                try:
                    results.append(self.process_file(file_path))
                except Exception as e:
                    print(f"Warning: Could not process {file_path}: {e}")
        return results

    def _process_pptx(self, file_path: Path) -> dict:
        """Extract content from PowerPoint."""
        prs = Presentation(file_path)
        slides = []
        all_text = []
        image_count = 0
        image_data: List[Dict[str, Any]] = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            slide_images = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image = shape.image
                        img_filename = (
                            f"slide_{slide_num}_img_{len(slide_images)}.{image.ext}"
                        )
                        slide_images.append(
                            {
                                "slide": slide_num,
                                "filename": img_filename,
                                "size": (shape.width, shape.height),
                            }
                        )
                        # Save image blob for Qwen3-VL embedding
                        image_data.append(
                            {
                                "slide": slide_num,
                                "blob": image.blob,
                                "ext": image.ext,
                                "filename": img_filename,
                            }
                        )
                        image_count += 1
                    except Exception:
                        pass

            text_content = "\n".join(slide_texts)
            if text_content:
                all_text.append(f"SLIDE {slide_num}:\n{text_content}")

            slides.append(
                {"slide_num": slide_num, "text": text_content, "images": slide_images}
            )

        return {
            "type": "pptx",
            "path": str(file_path),
            "slides": slides,
            "content": "\n\n".join(all_text),
            "image_count": image_count,
            "image_data": image_data,
            "chunks": [
                {"content": slide["text"], "type": "text", "slide": slide["slide_num"]}
                for slide in slides
                if slide["text"]
            ],
        }

    def _process_text(self, file_path: Path) -> dict:
        """Process text/markdown files."""
        content = file_path.read_text(encoding="utf-8")

        # Simple chunking by headers
        chunks = []
        lines = content.split("\n")
        current_chunk = []

        for line in lines:
            if line.strip().startswith("#") and current_chunk:
                chunks.append("\n".join(current_chunk))
                current_chunk = [line]
            else:
                current_chunk.append(line)

        if current_chunk:
            chunks.append("\n".join(current_chunk))

        return {
            "type": "text",
            "path": str(file_path),
            "content": content,
            "chunks": [
                {"content": chunk, "type": "text"}
                for chunk in chunks
                if chunk.strip()
            ],
        }

    def _process_pdf(self, file_path: Path) -> dict:
        """Extract text and images from PDF using PyMuPDF (fitz)."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            return {
                "type": "pdf",
                "path": str(file_path),
                "content": f"[PDF: {file_path.name} — PyMuPDF not installed. Install with: pip install pymupdf]",
                "chunks": [],
                "page_count": 0,
                "image_count": 0,
            }

        doc = fitz.open(file_path)
        all_text = []
        chunks = []
        image_data = []
        image_count = 0

        for page_num, page in enumerate(doc, 1):
            # Extract text
            text = page.get_text()
            if text.strip():
                all_text.append(f"PAGE {page_num}:\n{text.strip()}")
                chunks.append({
                    "content": text.strip(),
                    "type": "text",
                    "page": page_num,
                })

            # Extract images (figures, charts, diagrams)
            if not self.use_api:  # Only extract images for local processing
                try:
                    image_list = page.get_images(full=True)
                    for img_index, img in enumerate(image_list, 1):
                        try:
                            xref = img[0]
                            base_image = doc.extract_image(xref)
                            image_bytes = base_image["image"]
                            image_ext = base_image["ext"]
                            img_filename = f"page_{page_num}_img_{img_index}.{image_ext}"
                            
                            image_data.append({
                                "page": page_num,
                                "filename": img_filename,
                                "blob": image_bytes,
                                "ext": image_ext,
                            })
                            image_count += 1
                        except Exception:
                            pass  # Skip problematic images
                except Exception:
                    pass  # Skip image extraction for this page

        doc.close()

        return {
            "type": "pdf",
            "path": str(file_path),
            "content": "\n\n".join(all_text),
            "chunks": chunks,
            "page_count": len(chunks),
            "image_count": image_count,
            "image_data": image_data,
        }

    def _process_docx(self, file_path: Path) -> dict:
        """Extract text from DOCX using python-docx."""
        try:
            import docx
        except ImportError:
            return {
                "type": "docx",
                "path": str(file_path),
                "content": f"[DOCX: {file_path.name} — python-docx not installed. Install with: pip install python-docx]",
                "chunks": [],
                "paragraph_count": 0,
            }

        doc = docx.Document(file_path)
        paragraphs = []
        chunks = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
                chunks.append({"content": text, "type": "text"})

        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                if row_text:
                    table_row = " | ".join(row_text)
                    paragraphs.append(table_row)
                    chunks.append({"content": table_row, "type": "table_row"})

        return {
            "type": "docx",
            "path": str(file_path),
            "content": "\n\n".join(paragraphs),
            "chunks": chunks,
            "paragraph_count": len(paragraphs),
        }

    def _process_csv(self, file_path: Path) -> dict:
        """Process CSV files using standard library csv module."""
        import csv
        
        rows = []
        chunks = []
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                reader = csv.reader(f)
                headers = next(reader, None)
                if headers:
                    rows.append(headers)
                
                for row in reader:
                    if any(cell.strip() for cell in row):
                        rows.append(row)
                        # Create a natural language representation
                        row_text = " | ".join(f"{headers[i] if headers else f'Col{i}'}: {cell}" 
                                             for i, cell in enumerate(row) if cell.strip())
                        if row_text:
                            chunks.append({"content": row_text, "type": "csv_row"})
        except Exception as e:
            return {
                "type": "csv",
                "path": str(file_path),
                "content": f"[CSV: {file_path.name} — Error reading: {e}]",
                "chunks": [],
                "row_count": 0,
            }

        # Create header summary
        content_parts = []
        if headers:
            content_parts.append(f"Columns: {', '.join(headers)}")
        content_parts.append(f"Total rows: {len(rows) - (1 if headers else 0)}")
        
        # Add sample rows (first 10)
        for i, row in enumerate(rows[:10], 1):
            row_text = " | ".join(row)
            content_parts.append(f"Row {i}: {row_text}")

        return {
            "type": "csv",
            "path": str(file_path),
            "content": "\n".join(content_parts),
            "chunks": chunks,
            "row_count": len(rows),
            "headers": headers,
        }

    def _process_excel(self, file_path: Path) -> dict:
        """Process Excel files using openpyxl (for .xlsx) or xlrd (for .xls)."""
        try:
            import openpyxl
        except ImportError:
            return {
                "type": "excel",
                "path": str(file_path),
                "content": f"[Excel: {file_path.name} — openpyxl not installed. Install with: pip install openpyxl]",
                "chunks": [],
                "sheet_count": 0,
            }

        try:
            wb = openpyxl.load_workbook(file_path, data_only=True)
        except Exception as e:
            return {
                "type": "excel",
                "path": str(file_path),
                "content": f"[Excel: {file_path.name} — Error reading: {e}]",
                "chunks": [],
                "sheet_count": 0,
            }

        all_sheets = []
        chunks = []
        sheet_info = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_rows = []
            
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None and str(cell).strip() for cell in row):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    sheet_rows.append(row_values)
                    
                    # Create natural language chunk
                    row_text = " | ".join(cell for cell in row_values if cell.strip())
                    if row_text:
                        chunks.append({
                            "content": f"[Sheet: {sheet_name}] {row_text}",
                            "type": "excel_row",
                            "sheet": sheet_name,
                        })

            if sheet_rows:
                sheet_info.append(f"Sheet '{sheet_name}': {len(sheet_rows)} rows")
                all_sheets.append(f"\n=== Sheet: {sheet_name} ===")
                for i, row in enumerate(sheet_rows[:20], 1):  # First 20 rows per sheet
                    row_text = " | ".join(row)
                    all_sheets.append(f"Row {i}: {row_text}")
                if len(sheet_rows) > 20:
                    all_sheets.append(f"... ({len(sheet_rows) - 20} more rows)")

        wb.close()

        return {
            "type": "excel",
            "path": str(file_path),
            "content": "\n".join(sheet_info + all_sheets),
            "chunks": chunks,
            "sheet_count": len(wb.sheetnames),
            "sheet_names": wb.sheetnames,
        }

    # ── Media processing (audio / video) ──────────────────────────

    @staticmethod
    def _extract_audio_ffmpeg(media_path: str, output_path: str) -> bool:
        """Extract audio track from a media file using ffmpeg.

        Args:
            media_path: Path to video/audio file.
            output_path: Path for the extracted .wav file.

        Returns:
            True if extraction succeeded, False otherwise.
        """
        import subprocess

        cmd = [
            "ffmpeg", "-i", media_path,
            "-vn",                    # drop video
            "-acodec", "pcm_s16le",   # 16-bit PCM (Parakeet-compatible)
            "-ar", "16000",           # 16 kHz sample rate
            "-ac", "1",               # mono
            "-y",                     # overwrite
            "-loglevel", "error",
            output_path,
        ]
        try:
            subprocess.run(cmd, check=True, capture_output=True)
            return True
        except (subprocess.CalledProcessError, FileNotFoundError) as exc:
            print(f"  Warning: ffmpeg audio extraction failed: {exc}")
            return False

    @staticmethod
    def _extract_frames_ffmpeg(
        video_path: str, fps: float = 0.5, max_frames: int = 32
    ) -> List[str]:
        """Extract key frames from a video using ffmpeg.

        Returns:
            List of temporary frame image paths.
        """
        import subprocess
        import tempfile

        tmp_dir = tempfile.mkdtemp(prefix="vrlmrag_frames_")
        out_pattern = os.path.join(tmp_dir, "frame_%04d.jpg")

        cmd = [
            "ffmpeg", "-i", video_path,
            "-vf", f"fps={fps}",
            "-frames:v", str(max_frames),
            "-q:v", "2",
            "-loglevel", "error",
            out_pattern,
        ]
        try:
            subprocess.run(cmd, check=True, capture_output=True)
        except (subprocess.CalledProcessError, FileNotFoundError):
            return []

        frames = sorted(
            os.path.join(tmp_dir, f)
            for f in os.listdir(tmp_dir)
            if f.endswith(".jpg")
        )
        return frames

    def _process_media(self, file_path: Path) -> dict:
        """Process an audio or video file.

        Local mode (use_api=False):
            - Extracts audio track (ffmpeg) → transcribes with Parakeet
            - Extracts key frames for Qwen3-VL visual embedding

        API mode (use_api=True):
            - Extracts key frames → ZenMux omni model describes them
            - Extracts audio → ZenMux omni model transcribes/summarises
            - Falls back to frame descriptions if audio API unavailable

        Returns a document dict with text chunks (transcript) and
        frame_paths / audio_path metadata for the embedding stage.
        
        SAFETY: This function wraps all operations in try-except to ensure
        media processing never crashes the entire system.
        """
        import tempfile

        ext = file_path.suffix.lower()
        is_video = ext in _VIDEO_EXTENSIONS
        media_type = "video" if is_video else "audio"

        result: Dict[str, Any] = {
            "type": media_type,
            "path": str(file_path),
            "content": "",
            "chunks": [],
            "frame_paths": [],
            "audio_path": None,
        }

        try:
            # --- Extract audio track ---
            audio_path: Optional[str] = None
            if is_video:
                audio_path = tempfile.mktemp(suffix=".wav", prefix="vrlmrag_audio_")
                ok = self._extract_audio_ffmpeg(str(file_path), audio_path)
                if not ok or not Path(audio_path).exists() or Path(audio_path).stat().st_size < 100:
                    audio_path = None
            else:
                # Already an audio file
                audio_path = str(file_path)

            result["audio_path"] = audio_path

            # --- Transcribe audio (local Parakeet) ---
            transcript = ""
            if audio_path and self.transcription_provider is not None:
                try:
                    print(f"  [media] Transcribing audio with Parakeet...")
                    transcript = self.transcription_provider.transcribe(audio_path)
                    if isinstance(transcript, dict):
                        transcript = transcript.get("text", "")
                    print(f"  [media] Transcript: {len(transcript)} chars")
                except Exception as exc:
                    print(f"  Warning: Parakeet transcription failed: {exc}")

            # --- Extract key frames (video only) ---
            frame_paths: List[str] = []
            if is_video:
                print(f"  [media] Extracting key frames...")
                frame_paths = self._extract_frames_ffmpeg(str(file_path))
                print(f"  [media] Extracted {len(frame_paths)} frames")
                result["frame_paths"] = frame_paths

            # --- Build content and chunks ---
            content_parts: List[str] = []
            chunks: List[dict] = []

            if transcript:
                content_parts.append(transcript)
                # Chunk transcript into ~500-char segments for embedding
                words = transcript.split()
                chunk_words: List[str] = []
                for word in words:
                    chunk_words.append(word)
                    if len(" ".join(chunk_words)) >= 500:
                        chunk_text = " ".join(chunk_words)
                        chunks.append({"content": chunk_text, "type": "transcript"})
                        chunk_words = []
                if chunk_words:
                    chunk_text = " ".join(chunk_words)
                    if chunk_text.strip():
                        chunks.append({"content": chunk_text, "type": "transcript"})

            if not content_parts:
                # No transcript available — placeholder
                content_parts.append(f"[{media_type.title()}: {file_path.name}]")

            result["content"] = "\n\n".join(content_parts)
            result["chunks"] = chunks

        except Exception as e:
            # CRITICAL SAFETY: Never let media processing crash the system
            print(f"  [CRITICAL] Media processing failed for {file_path.name}: {e}")
            print(f"  [CRITICAL] Returning empty document to prevent system crash.")
            result["content"] = f"[{media_type.title()}: {file_path.name} — processing failed]"
            result["chunks"] = []
            result["frame_paths"] = []
            result["audio_path"] = None

        return result


class MarkdownReportGenerator:
    """Generate markdown reports from VL-RAG-Graph-RLM results."""

    def generate(self, results: dict, output_path: Optional[str] = None) -> str:
        """Generate a markdown report."""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        lines = [
            "# VL-RAG-Graph-RLM Analysis Report",
            "",
            f"**Generated:** {timestamp}",
            f"**Provider:** {results.get('provider', 'N/A')}",
            f"**Model:** {results.get('model', 'N/A')}",
            f"**Embeddings:** {results.get('embedding_model', 'N/A')}",
            f"**Reranker:** {results.get('reranker_model', 'N/A')}",
            "",
            "## Summary",
            "",
            f"- **Documents Processed:** {results.get('document_count', 0)}",
            f"- **Total Chunks:** {results.get('total_chunks', 0)}",
            f"- **Embedded Documents:** {results.get('embedded_count', 0)}",
            f"- **Processing Time:** {results.get('execution_time', 0):.2f}s",
            "",
            "## Knowledge Graph",
            "",
            results.get("knowledge_graph", "No knowledge graph generated"),
            "",
            "## Query Responses",
            "",
        ]

        for query_response in results.get("queries", []):
            lines.extend(
                [
                    f"### Query: {query_response['query']}",
                    "",
                    query_response["response"],
                    "",
                    f"*Time: {query_response.get('time', 0):.2f}s*",
                    "",
                ]
            )

            # Add source info if available
            if query_response.get("sources"):
                lines.append("**Retrieved Sources:**")
                for src in query_response["sources"][:5]:
                    score = src.get("score", 0)
                    content_preview = src.get("content", "")[:100]
                    lines.append(f"- [Score: {score:.2f}] {content_preview}...")
                lines.append("")

        lines.extend(["## Source Documents", ""])

        for doc in results.get("documents", []):
            lines.extend(
                [
                    f"### {Path(doc['path']).name}",
                    "",
                    f"- **Type:** {doc['type']}",
                    f"- **Path:** `{doc['path']}`",
                ]
            )
            if "slide_count" in doc:
                lines.append(f"- **Slides:** {doc['slide_count']}")
            if "image_count" in doc:
                lines.append(f"- **Images:** {doc['image_count']}")
            lines.append("")

        markdown = "\n".join(lines)

        if output_path:
            Path(output_path).write_text(markdown, encoding="utf-8")
            print(f"\nReport saved to: {output_path}")

        return markdown


SUPPORTED_PROVIDERS: Dict[str, Dict[str, Any]] = {
    "sambanova": {
        "env_key": "SAMBANOVA_API_KEY",
        "url": "https://cloud.sambanova.ai",
        "description": "SambaNova Cloud — DeepSeek-V3-0324 (200+ tok/s, 32K context, 200K TPD free)",
        "context_budget": 32000,
    },
    "nebius": {
        "env_key": "NEBIUS_API_KEY",
        "url": "https://tokenfactory.nebius.com",
        "description": "Nebius Token Factory — MiniMax-M2.1 (128K context, no daily limits)",
        "context_budget": 100000,
    },
    "openrouter": {
        "env_key": "OPENROUTER_API_KEY",
        "url": "https://openrouter.ai",
        "description": "OpenRouter — 200+ models, pay-per-token routing",
        "context_budget": 32000,
    },
    "openai": {
        "env_key": "OPENAI_API_KEY",
        "url": "https://platform.openai.com",
        "description": "OpenAI — GPT-4o-mini (128K context)",
        "context_budget": 32000,
    },
    "anthropic": {
        "env_key": "ANTHROPIC_API_KEY",
        "url": "https://console.anthropic.com",
        "description": "Anthropic — Claude 3.5 Haiku (200K context)",
        "context_budget": 32000,
    },
    "gemini": {
        "env_key": "GOOGLE_API_KEY",
        "url": "https://makersuite.google.com/app/apikey",
        "description": "Google Gemini — gemini-1.5-flash (1M context)",
        "context_budget": 64000,
    },
    "groq": {
        "env_key": "GROQ_API_KEY",
        "url": "https://console.groq.com",
        "description": "Groq — Ultra-fast LPU inference (Kimi K2, GPT-OSS-120B, Llama 4)",
        "context_budget": 32000,
    },
    "deepseek": {
        "env_key": "DEEPSEEK_API_KEY",
        "url": "https://platform.deepseek.com",
        "description": "DeepSeek — DeepSeek-V3 / R1 reasoning",
        "context_budget": 32000,
    },
    "mistral": {
        "env_key": "MISTRAL_API_KEY",
        "url": "https://console.mistral.ai",
        "description": "Mistral AI — mistral-large-latest (128K context)",
        "context_budget": 32000,
    },
    "fireworks": {
        "env_key": "FIREWORKS_API_KEY",
        "url": "https://fireworks.ai",
        "description": "Fireworks AI — Serverless open-source models",
        "context_budget": 32000,
    },
    "together": {
        "env_key": "TOGETHER_API_KEY",
        "url": "https://api.together.ai",
        "description": "Together AI — Open-source model hosting",
        "context_budget": 32000,
    },
    "zenmux": {
        "env_key": "ZENMUX_API_KEY",
        "url": "https://zenmux.ai",
        "description": "ZenMux — Unified API gateway (59+ models, OpenAI/Anthropic protocols)",
        "context_budget": 32000,
    },
    "zai": {
        "env_key": "ZAI_API_KEY",
        "url": "https://open.bigmodel.cn",
        "description": "z.ai (Zhipu AI) — GLM-4.7 (Coding Plan first, then normal API)",
        "context_budget": 32000,
    },
    "azure_openai": {
        "env_key": "AZURE_OPENAI_API_KEY",
        "url": "https://portal.azure.com",
        "description": "Azure OpenAI — Enterprise GPT deployments",
        "context_budget": 32000,
    },
    "cerebras": {
        "env_key": "CEREBRAS_API_KEY",
        "url": "https://cloud.cerebras.ai",
        "description": "Cerebras — Ultra-fast wafer-scale (GLM-4.7, GPT-OSS-120B, Qwen3-235B)",
        "context_budget": 32000,
    },
    "ollama": {
        "env_key": "OLLAMA_MODEL",  # Not really an API key, but indicates intent to use Ollama
        "url": "https://ollama.com",
        "description": "Ollama — Local LLM inference (llama3.2, mistral, qwen2.5, etc.)",
        "context_budget": 32000,
    },
}


def list_providers() -> None:
    """Print all supported providers with their status."""
    print(f"vrlmrag v{__version__} — Supported Providers\n")
    print(f"{'Provider':<16} {'API Key':<12} {'Description'}")
    print("-" * 80)
    for name, info in SUPPORTED_PROVIDERS.items():
        key_set = "✓ SET" if os.getenv(info["env_key"]) else "✗ unset"
        print(f"{name:<16} {key_set:<12} {info['description']}")
    print()
    print("Set API keys in .env or via environment variables.")
    print("See .env.example for all configuration options.")


def show_hierarchy() -> None:
    """Print the provider hierarchy with availability status."""
    hierarchy = get_hierarchy()
    available = get_available_providers(hierarchy)

    print(f"vrlmrag v{__version__} — Provider Hierarchy\n")
    print("Providers are tried in order. First available provider is used.")
    print("Edit PROVIDER_HIERARCHY in .env to customize the order.\n")
    print(f"{'#':<4} {'Provider':<16} {'Status':<12} {'API Key Env Var'}")
    print("-" * 60)
    for i, provider in enumerate(hierarchy, 1):
        env_key = PROVIDER_KEY_MAP.get(provider, f"{provider.upper()}_API_KEY")
        if provider in available:
            status = "✓ READY"
        else:
            status = "✗ no key"
        print(f"{i:<4} {provider:<16} {status:<12} {env_key}")

    print(f"\nAvailable: {len(available)}/{len(hierarchy)} providers ready")
    if available:
        print(f"Auto mode will use: {available[0]}")
    print(f"\nCustomize: PROVIDER_HIERARCHY={','.join(hierarchy)}")


def run_analysis(
    provider: str,
    input_path: str,
    query: Optional[str] = None,
    output: Optional[str] = None,
    model: Optional[str] = None,
    max_depth: int = 3,
    max_iterations: int = 10,
    use_api: bool = False,
    text_only: bool = False,
    _quiet: bool = False,
    multi_query: bool = False,
    rrf_weights: Optional[List[float]] = None,
    use_graph_augmented: bool = False,
    graph_hops: int = 2,
    output_format: str = "markdown",
    verbose: bool = False,
) -> dict:
    """Run full VL-RAG-Graph-RLM analysis with any supported provider.

    Pipeline (6 pillars):
      1. Document intake (PPTX, TXT, MD)
      2. Qwen3-VL embedding (text + images → unified vector space)
      3. Hybrid search (dense + keyword + RRF fusion)
      4. Qwen3-VL reranking (cross-attention relevance scoring)
      5. Knowledge graph extraction via RLM
      6. Query answering via RLM with retrieved context → markdown report
    """
    # Helper for conditional printing
    def _print(*args, **kwargs):
        if not _quiet:
            print(*args, **kwargs)

    # Handle 'auto' provider — resolve from hierarchy
    is_auto = provider == "auto"
    if is_auto:
        try:
            provider = resolve_auto_provider()
            _print(f"[auto] Resolved provider from hierarchy: {provider}")
        except RuntimeError as e:
            _print(f"Error: {e}")
            _print("Run 'vrlmrag --show-hierarchy' to see provider status.")
            sys.exit(1)

    if provider not in SUPPORTED_PROVIDERS:
        _print(f"Error: Unknown provider '{provider}'")
        _print(f"Supported: {', '.join(SUPPORTED_PROVIDERS.keys())}")
        sys.exit(1)

    prov_info = SUPPORTED_PROVIDERS[provider]
    context_budget = prov_info["context_budget"]

    # Check API key
    api_key = os.getenv(prov_info["env_key"])
    if not api_key:
        _print(f"Error: {prov_info['env_key']} not set")
        _print(f"Get your API key from: {prov_info['url']}")
        sys.exit(1)

    # Resolve model (env var → explicit arg → provider default)
    resolved_model = model  # explicit --model flag takes priority
    if not resolved_model:
        resolved_model = None  # let VLRAGGraphRLM resolve from env/defaults

    _embed_mode = "API" if use_api else ("text-only" if text_only else "local Qwen3-VL")
    _print("=" * 70)
    _print(f"vrlmrag v{__version__} — Full VL-RAG-Graph-RLM Pipeline")
    _print("=" * 70)
    _print(f"Provider:       {provider}")
    _print(f"Model:          {resolved_model or '(auto-detect from env/defaults)'}")
    _print(f"Embedding:      {_embed_mode}")
    _print(f"Context budget: {context_budget:,} chars")
    _print(f"Input:          {input_path}")
    _print("=" * 70)

    overall_start = time.time()

    # ================================================================
    # Persistence directory — used for embeddings + knowledge graph
    # Created early so manifest-based change detection can work.
    # ================================================================
    _store_path = Path(input_path)
    if _store_path.is_file():
        store_dir = _store_path.parent / ".vrlmrag_store"
    else:
        store_dir = _store_path / ".vrlmrag_store"
    store_dir.mkdir(parents=True, exist_ok=True)

    # ================================================================
    # Manifest-based change detection — skip unchanged files
    # ================================================================
    _current_files = _scan_supported_files(Path(input_path))
    _manifest = _load_manifest(store_dir)
    _new_or_modified, _deleted = _detect_file_changes(_current_files, _manifest)
    _embeddings_exist = (store_dir / "embeddings.json").exists() or (store_dir / "embeddings_text.json").exists()
    _needs_processing = bool(_new_or_modified) or not _embeddings_exist

    if _embeddings_exist and not _needs_processing:
        _print(f"\n[1/6] Store up-to-date — {len(_current_files)} file(s) unchanged, skipping document processing")
    else:
        if _embeddings_exist and _new_or_modified:
            _print(f"\n[1/6] Incremental update — {len(_new_or_modified)} new/modified file(s)")
        else:
            _print(f"\n[1/6] Processing documents: {input_path}")

    # ================================================================
    # Step 1: Document Processing
    # Local mode: Parakeet ASR transcribes audio → text chunks
    # API mode: audio/video handled by omni model at embedding time
    # ================================================================
    documents: List[dict] = []
    all_chunks: List[dict] = []

    if _needs_processing:
        _transcriber = None
        if not use_api and not text_only and HAS_PARAKEET:
            _print("  [init] Parakeet ASR available for audio transcription")
            _transcriber = create_parakeet_transcriber(
                cache_dir=str(Path.home() / ".vrlmrag" / "parakeet_cache"),
            )
        processor = DocumentProcessor(
            transcription_provider=_transcriber,
            use_api=use_api,
        )

        if _new_or_modified and _embeddings_exist:
            # Incremental: only process changed files
            for fpath in _new_or_modified:
                result = processor.process_path(fpath)
                if isinstance(result, dict):
                    documents.append(result)
                elif isinstance(result, list):
                    documents.extend(result)
        else:
            # Full processing (first run or no store)
            result = processor.process_path(input_path)
            if isinstance(result, dict):
                documents = [result]
            elif isinstance(result, list):
                documents = result

        for doc in documents:
            all_chunks.extend(doc.get("chunks", []))

        _print(f"  Processed {len(documents)} document(s), {len(all_chunks)} chunks")

    # ================================================================
    # Step 2: Qwen3-VL Embedding + Vector Store
    # Local models acquire the cross-process lock to ensure only one
    # model is loaded in RAM at a time (across all CLI/MCP sessions).
    # API-based embeddings bypass the lock entirely.
    # ================================================================
    embedding_model_name = os.getenv("VRLMRAG_LOCAL_EMBEDDING_MODEL", "Qwen/Qwen3-VL-Embedding-2B")
    reranker_model_name = os.getenv("VRLMRAG_RERANKER_MODEL", "ms-marco-MiniLM-L-12-v2")
    embedded_count = 0

    store = None
    reranker_vl = None
    _lock_ctx = None

    if text_only and HAS_TEXT_EMBEDDING:
        text_model = os.getenv("VRLMRAG_TEXT_ONLY_MODEL", "Qwen/Qwen3-Embedding-0.6B")
        _lock_ctx = local_model_lock(text_model, description="CLI run_analysis (text-only)")
        _lock_ctx.__enter__()
        _print(f"\n[2/6] Loading text-only embedding ({text_model})...")
        embedder = create_text_embedder(model_name=text_model)
        embedding_model_name = text_model
        _print(f"  Embedding dim: {embedder.embedding_dim}")

        storage_file = str(store_dir / "embeddings_text.json")

        store = MultimodalVectorStore(
            embedding_provider=embedder,
            storage_path=storage_file,
        )
    elif use_api and HAS_API_EMBEDDING:
        _print("\n[2/6] Loading API-based embedding (OpenRouter + ZenMux omni)...")
        try:
            embedder = create_api_embedder()
            embedding_model_name = os.getenv("VRLMRAG_EMBEDDING_MODEL", "openai/text-embedding-3-small")
            vlm_model = os.getenv("VRLMRAG_VLM_MODEL", "inclusionai/ming-flash-omni-preview")
            _print(f"  Embedding model: {embedding_model_name}")
            _print(f"  Omni VLM model:  {vlm_model}")
        except Exception as e:
            _print(f"\n[ERROR] API embedding failed: {e}")
            if HAS_QWEN3VL:
                _print(f"\n[recovery] Falling back to local Qwen3-VL embeddings...")
                _print(f"[recovery] To use offline mode explicitly, run with --offline flag.")
                # Switch to local mode
                use_api = False
                text_only = False
                _lock_ctx = local_model_lock(embedding_model_name, description="CLI run_analysis (Qwen3-VL fallback)")
                _lock_ctx.__enter__()
                import torch
                device = "mps" if torch.backends.mps.is_available() else "cpu"
                _print(f"\n[2/6] Loading Qwen3-VL Embedding model ({device})...")
                embedder = create_qwen3vl_embedder(
                    model_name=embedding_model_name, device=device
                )
                _print(f"  Embedding dim: {embedder.embedding_dim}")
            else:
                _print(f"\n[FATAL] Cannot proceed — no embedding provider available.")
                _print(f"        API embedding failed and local Qwen3-VL is not installed.")
                _print(f"\n        To fix this:")
                _print(f"          1. Set valid OPENROUTER_API_KEY in .env for API mode")
                _print(f"          2. Or install local models: pip install torch transformers qwen-vl-utils")
                sys.exit(1)

        storage_file = str(store_dir / "embeddings.json")

        store = MultimodalVectorStore(
            embedding_provider=embedder,
            storage_path=storage_file,
        )
        existing_count = len(store.documents)
        if existing_count:
            _print(f"  Loaded {existing_count} existing embeddings from store")

        # Embed text chunks via API
        skipped_count = 0
        _print(f"\n[3/6] Embedding {len(all_chunks)} chunks via API...")
        for chunk in all_chunks:
            content = chunk.get("content", "")
            if not content.strip():
                continue
            if store.content_exists(content):
                skipped_count += 1
                continue
            metadata = {"type": chunk.get("type", "text")}
            if "slide" in chunk:
                metadata["slide"] = chunk["slide"]
            store.add_text(
                content=content, metadata=metadata,
                instruction=_DOCUMENT_INSTRUCTION,
            )
            embedded_count += 1

        # Embed images (PPTX slides) via omni model
        for doc in documents:
            for img_info in doc.get("image_data", []):
                try:
                    temp_path = f"/tmp/vrlmrag_{img_info['filename']}"
                    with open(temp_path, "wb") as f:
                        f.write(img_info["blob"])
                    prev_count = len(store.documents)
                    store.add_image(
                        image_path=temp_path,
                        description=f"Image from slide {img_info['slide']}",
                        metadata={
                            "type": "image",
                            "slide": img_info["slide"],
                            "filename": img_info["filename"],
                        },
                        instruction=_DOCUMENT_INSTRUCTION,
                    )
                    if len(store.documents) > prev_count:
                        embedded_count += 1
                    else:
                        skipped_count += 1
                except Exception as e:
                    print(f"  Warning: Could not embed image: {e}")

        # Embed video frames via omni model (ZenMux describes → embeds)
        for doc in documents:
            frame_paths = doc.get("frame_paths", [])
            if frame_paths:
                print(f"  Embedding {len(frame_paths)} video frames via omni model...")
                for i, frame_path in enumerate(frame_paths):
                    try:
                        prev_count = len(store.documents)
                        store.add_image(
                            image_path=frame_path,
                            description=f"Video frame {i+1} from {Path(doc['path']).name}",
                            metadata={
                                "type": "video_frame",
                                "frame_index": i,
                                "source": doc["path"],
                            },
                            instruction=_DOCUMENT_INSTRUCTION,
                        )
                        if len(store.documents) > prev_count:
                            embedded_count += 1
                        else:
                            skipped_count += 1
                    except Exception as e:
                        print(f"  Warning: Could not embed frame {i}: {e}")

        print(f"  New embeddings: {embedded_count} | Skipped: {skipped_count}")
        print(f"  Total in store: {len(store.documents)} documents")
        print(f"  Store persisted to: {storage_file}")

        # Update manifest after successful indexing
        if _needs_processing:
            _manifest["files"] = {str(k): v for k, v in _current_files.items()}
            _save_manifest(store_dir, _manifest)

        # FlashRank reranker (lightweight, works with API mode too)
        if HAS_FLASHRANK:
            _print("\n[4/6] Loading FlashRank Reranker...")
            reranker_vl = create_flashrank_reranker(model_name=reranker_model_name)
            _print(f"  Reranker loaded ({reranker_model_name})")
    elif HAS_QWEN3VL:
        _lock_ctx = local_model_lock(embedding_model_name, description="CLI run_analysis (Qwen3-VL)")
        _lock_ctx.__enter__()

        import torch

        device = "mps" if torch.backends.mps.is_available() else "cpu"
        _print(f"\n[2/6] Loading Qwen3-VL Embedding model ({device})...")
        embedder = create_qwen3vl_embedder(
            model_name=embedding_model_name, device=device
        )
        _print(f"  Embedding dim: {embedder.embedding_dim}")

        storage_file = str(store_dir / "embeddings.json")

        store = MultimodalVectorStore(
            embedding_provider=embedder,
            storage_path=storage_file,
        )
        existing_count = len(store.documents)
        if existing_count:
            _print(f"  Loaded {existing_count} existing embeddings from store")

        # Embed text chunks (dedup: skips already-embedded content)
        skipped_count = 0
        _print(f"\n[3/6] Embedding {len(all_chunks)} chunks with Qwen3-VL...")
        for i, chunk in enumerate(all_chunks):
            content = chunk.get("content", "")
            if not content.strip():
                continue
            if store.content_exists(content):
                skipped_count += 1
                continue
            metadata = {"type": chunk.get("type", "text")}
            if "slide" in chunk:
                metadata["slide"] = chunk["slide"]
            store.add_text(
                content=content, metadata=metadata,
                instruction=_DOCUMENT_INSTRUCTION,
            )
            embedded_count += 1
            if (embedded_count) % 5 == 0:
                _print(f"  Embedded {embedded_count} new chunks so far...")

        # Embed any extracted images (dedup handled inside store)
        for doc in documents:
            for img_info in doc.get("image_data", []):
                try:
                    temp_path = f"/tmp/vrlmrag_{img_info['filename']}"
                    with open(temp_path, "wb") as f:
                        f.write(img_info["blob"])
                    prev_count = len(store.documents)
                    store.add_image(
                        image_path=temp_path,
                        description=f"Image from slide {img_info['slide']}",
                        metadata={
                            "type": "image",
                            "slide": img_info["slide"],
                            "filename": img_info["filename"],
                        },
                        instruction=_DOCUMENT_INSTRUCTION,
                    )
                    if len(store.documents) > prev_count:
                        embedded_count += 1
                    else:
                        skipped_count += 1
                except Exception as e:
                    _print(f"  Warning: Could not embed image {img_info['filename']}: {e}")

        # Embed video frames from media documents
        for doc in documents:
            frame_paths = doc.get("frame_paths", [])
            if frame_paths:
                _print(f"  Embedding {len(frame_paths)} video frames...")
                for i, frame_path in enumerate(frame_paths):
                    try:
                        prev_count = len(store.documents)
                        store.add_image(
                            image_path=frame_path,
                            description=f"Video frame {i+1} from {Path(doc['path']).name}",
                            metadata={
                                "type": "video_frame",
                                "frame_index": i,
                                "source": doc["path"],
                            },
                            instruction=_DOCUMENT_INSTRUCTION,
                        )
                        if len(store.documents) > prev_count:
                            embedded_count += 1
                        else:
                            skipped_count += 1
                    except Exception as e:
                        _print(f"  Warning: Could not embed frame {i}: {e}")

        _print(f"  New embeddings: {embedded_count} | Skipped (already in store): {skipped_count}")
        _print(f"  Total in store: {len(store.documents)} documents")
        _print(f"  Store persisted to: {storage_file}")

        # Load lightweight FlashRank reranker (~34 MB — coexists with embedder)
        if HAS_FLASHRANK:
            _print("\n[4/6] Loading FlashRank Reranker...")
            reranker_vl = create_flashrank_reranker(model_name=reranker_model_name)
            _print(f"  Reranker loaded ({reranker_model_name})")
        else:
            _print("\n[4/6] FlashRank not available — using RRF-only retrieval")
            _print("  Install with: pip install flashrank")
    else:
        _print("\n[2/6] Qwen3-VL not available — using fallback text reranking")
        _print("  Install with: pip install torch transformers qwen-vl-utils torchvision")
        embedding_model_name = "N/A (fallback)"
        reranker_model_name = "N/A (fallback)"

    # Release the local model lock after embedding is complete.
    # RLM calls below are API-based and don't need the lock.
    if _lock_ctx is not None:
        _lock_ctx.__exit__(None, None, None)
        _lock_ctx = None

    # ================================================================
    # Step 3: Initialize RLM
    # ================================================================
    rlm_kwargs: Dict[str, Any] = {
        "provider": provider,
        "temperature": 0.0,
        "max_depth": max_depth,
        "max_iterations": max_iterations,
    }
    if resolved_model:
        rlm_kwargs["model"] = resolved_model

    _print(f"\n[5/6] Initializing RLM ({provider})...")
    rlm = VLRAGGraphRLM(**rlm_kwargs)
    _print(f"  Model: {rlm.model}")

    # If auto mode, set up fallback hierarchy for query errors
    fallback_hierarchy = get_available_providers() if is_auto else None

    # ================================================================
    # Step 4: Build / Extend Knowledge Graph (persistent)
    # ================================================================
    kg_file = store_dir / "knowledge_graph.md"
    knowledge_graph = _load_knowledge_graph(kg_file)
    if knowledge_graph:
        _print(f"\n  Loaded existing knowledge graph ({len(knowledge_graph):,} chars)")

    _print("\n  Building knowledge graph for new content...")
    kg_char_limit = min(context_budget, 25000)
    kg_doc_limit = max(2000, kg_char_limit // max(len(documents), 1))
    kg_context = "\n\n".join([d.get("content", "")[:kg_doc_limit] for d in documents])
    try:
        kg_result = rlm.completion(
            _KG_EXTRACTION_PROMPT,
            kg_context[:kg_char_limit],
        )
        knowledge_graph = _merge_knowledge_graphs(knowledge_graph, kg_result.response)
        _save_knowledge_graph(kg_file, knowledge_graph)
        _print(f"  Knowledge graph persisted ({len(knowledge_graph):,} chars)")
    except Exception as e:
        if not knowledge_graph:
            knowledge_graph = f"Could not build knowledge graph: {e}"
        _print(f"  Warning: KG extraction failed: {e}")

    # ================================================================
    # Step 5: Run Queries with Qwen3-VL retrieval
    # ================================================================
    queries_to_run: List[str] = []
    if query:
        if multi_query:
            # Generate sub-queries for broader recall
            _print(f"[multi-query] Generating sub-queries for: {query}")
            sub_queries = generate_sub_queries(query, rlm)
            queries_to_run = sub_queries
            _print(f"[multi-query] Generated {len(sub_queries)} queries total")
        else:
            queries_to_run = [query]
    else:
        queries_to_run = [
            "What are the main topics covered?",
            "Summarize the key concepts presented.",
        ]

    _print(f"\n[6/6] Running {len(queries_to_run)} queries with full VL-RAG pipeline...")
    query_results: List[Dict[str, Any]] = []
    rrf = ReciprocalRankFusion(k=60)
    fallback_reranker = CompositeReranker()

    for q in queries_to_run:
        _print(f"\n  Query: {q}")
        qr = _run_vl_rag_query(
            q,
            store=store,
            reranker_vl=reranker_vl,
            rrf=rrf,
            fallback_reranker=fallback_reranker,
            all_chunks=all_chunks,
            knowledge_graph=knowledge_graph,
            context_budget=context_budget,
            rlm=rlm,
            fallback_hierarchy=fallback_hierarchy,
            provider=provider,
            resolved_model=resolved_model,
            max_depth=max_depth,
            max_iterations=max_iterations,
            verbose=not _quiet,
            rrf_weights=rrf_weights or [4.0, 1.0],
            use_graph_augmented=use_graph_augmented,
            graph_hops=graph_hops,
        )
        query_results.append(qr)

    # ================================================================
    # Compile Results & Generate Report
    # ================================================================
    total_elapsed = time.time() - overall_start
    results = {
        "provider": provider,
        "model": rlm.model,
        "embedding_model": embedding_model_name,
        "reranker_model": reranker_model_name,
        "document_count": len(documents),
        "total_chunks": len(all_chunks),
        "embedded_count": embedded_count,
        "execution_time": total_elapsed,
        "documents": documents,
        "knowledge_graph": knowledge_graph,
        "queries": query_results,
    }

    if not _quiet:
        if output_format == "json":
            # JSON output
            _print("\n" + json.dumps(results, indent=2, default=str))
        else:
            # Markdown report
            _print("\n" + "=" * 70)
            _print("Generating markdown report...")
            _print("=" * 70)

            generator = MarkdownReportGenerator()
            report = generator.generate(results, output)

            if not output:
                _print("\n" + report)

    return results


def _run_quality_check(
    collection_name: str,
    provider: str,
    model: Optional[str] = None,
    max_depth: int = 3,
    max_iterations: int = 10,
) -> None:
    """Run RLM-powered embedding quality assessment for a collection."""
    from vl_rag_graph_rlm.rlm_core import VLRAGGraphRLM
    from vl_rag_graph_rlm.collections import load_collection_meta, collection_load_kg
    
    print(f"\n[quality-check:{collection_name}] RLM-powered embedding quality assessment")
    print("=" * 60)
    
    meta = load_collection_meta(collection_name)
    slug = meta["name"]
    doc_count = meta.get("document_count", 0)
    chunk_count = meta.get("chunk_count", 0)
    embedding_model = meta.get("embedding_model", "unknown")
    
    print(f"Collection: {meta.get('display_name', slug)}")
    print(f"Documents: {doc_count} | Chunks: {chunk_count}")
    print(f"Embedding model: {embedding_model}")
    
    kg = collection_load_kg(slug)
    if not kg:
        print("\n[Warning] No knowledge graph found")
        kg_context = "No knowledge graph available."
    else:
        kg_context = kg[:8000]
        print(f"Knowledge graph: {len(kg):,} chars")
    
    rlm_kwargs = {
        "provider": provider,
        "temperature": 0.0,
        "max_depth": max_depth,
        "max_iterations": max_iterations,
    }
    if model:
        rlm_kwargs["model"] = model
    
    rlm = VLRAGGraphRLM(**rlm_kwargs)
    print(f"RLM: {rlm.model} (provider: {provider})")
    print("=" * 60)
    
    quality_prompt = """You are an embedding quality assessment system. Analyze the collection metadata and knowledge graph to evaluate retrieval quality.

Provide a quality assessment with:
- Overall Quality Score (0-100)
- Strengths: What's working well
- Weaknesses: What could be improved  
- Recommendations: Specific actions to improve quality

Respond with a structured assessment."""

    print("\nRunning RLM quality assessment...")
    try:
        assessment_context = f"""Collection: {meta.get('display_name', slug)}
Documents: {doc_count} | Chunks: {chunk_count}
Embedding Model: {embedding_model}
Reranker: {meta.get('reranker_model', 'unknown')}
Sources: {len(meta.get('sources', []))}

Knowledge Graph (first 8000 chars):
{kg_context}"""
        
        result = rlm.completion(quality_prompt, assessment_context[:12000])
        
        print("\n" + "=" * 60)
        print("QUALITY ASSESSMENT REPORT")
        print("=" * 60)
        print(result.response)
        print("=" * 60)
        
        import re
        score_match = re.search(r'(\d{1,3})\s*/\s*100|score[:\s]+(\d{1,3})', result.response, re.IGNORECASE)
        if score_match:
            score = score_match.group(1) or score_match.group(2)
            print(f"\nDetected Quality Score: {score}/100")
            if int(score) >= 80:
                print("Status: Excellent")
            elif int(score) >= 60:
                print("Status: Good")
            elif int(score) >= 40:
                print("Status: Fair — Reindexing may help")
            else:
                print("Status: Poor — Reindexing recommended")
    except Exception as e:
        print(f"\n[Error] Quality assessment failed: {e}")


# Backward-compatible aliases
def run_sambanova_analysis(
    input_path: str, query: Optional[str] = None, output: Optional[str] = None
) -> dict:
    """Backward-compatible alias for run_analysis(provider='sambanova', ...)."""
    return run_analysis("sambanova", input_path, query=query, output=output)


def run_nebius_analysis(
    input_path: str, query: Optional[str] = None, output: Optional[str] = None
) -> dict:
    """Backward-compatible alias for run_analysis(provider='nebius', ...)."""
    return run_analysis("nebius", input_path, query=query, output=output)


def _try_fallback_query(
    query: str,
    context: str,
    hierarchy: List[str],
    failed_provider: str,
    model: Optional[str],
    max_depth: int,
    max_iterations: int,
) -> Optional[Dict[str, Any]]:
    """Try fallback providers for a failed query.

    Note: model is NOT forwarded to fallback providers because model names
    are provider-specific (e.g., 'DeepSeek-V3.2' on SambaNova vs 'glm-4.7'
    on z.ai). Each fallback provider uses its own default model from env or
    hardcoded defaults.
    """
    for fb_provider in hierarchy:
        if fb_provider == failed_provider:
            continue
        try:
            fb_kwargs: Dict[str, Any] = {
                "provider": fb_provider,
                "temperature": 0.0,
                "max_depth": max_depth,
                "max_iterations": max_iterations,
            }

            fb_rlm = VLRAGGraphRLM(**fb_kwargs)
            q_start = time.time()
            result = fb_rlm.completion(query, context)
            elapsed = time.time() - q_start

            return {
                "query": query,
                "response": result.response,
                "time": elapsed,
                "sources": [],
                "fallback_provider": fb_provider,
            }
        except Exception:
            continue
    return None


def _keyword_search(
    store: "MultimodalVectorStore", query: str, top_k: int = 20
) -> List[SearchResult]:
    """BM25 keyword search across the vector store documents.
    
    Uses BM25 ranking algorithm (state-of-the-art for keyword retrieval).
    Falls back to simple token overlap if rank-bm25 not installed.
    """
    import re
    import math
    
    query_terms = set(re.findall(r"\w+", query.lower()))
    if not query_terms:
        return []
    
    # Try to use rank-bm25 if available
    try:
        from rank_bm25 import BM25Okapi
        
        # Prepare corpus
        corpus = []
        doc_ids = []
        for doc_id, doc in store.documents.items():
            tokens = re.findall(r"\w+", doc.content.lower())
            corpus.append(tokens)
            doc_ids.append(doc_id)
        
        if not corpus:
            return []
        
        # Build BM25 index
        bm25 = BM25Okapi(corpus)
        
        # Score query
        query_tokens = list(query_terms)
        scores = bm25.get_scores(query_tokens)
        
        # Create results
        results = []
        for idx, score in enumerate(scores):
            if score > 0:
                doc_id = doc_ids[idx]
                doc = store.get(doc_id)
                if doc:
                    results.append(
                        SearchResult(
                            id=doc_id,
                            content=doc.content,
                            metadata=doc.metadata,
                            keyword_score=score,
                            composite_score=score,
                        )
                    )
        
        results.sort(key=lambda x: x.keyword_score, reverse=True)
        return results[:top_k]
        
    except ImportError:
        # Fallback to simple token overlap if rank-bm25 not installed
        results = []
        for doc in store.documents.values():
            content_lower = doc.content.lower()
            matches = sum(1 for term in query_terms if term in content_lower)
            
            if matches > 0:
                score = matches / len(query_terms)
                results.append(
                    SearchResult(
                        id=doc.id,
                        content=doc.content,
                        metadata=doc.metadata,
                        keyword_score=score,
                        composite_score=score,
                    )
                )
        
        results.sort(key=lambda x: x.keyword_score, reverse=True)
        return results[:top_k]


def generate_sub_queries(original_query: str, rlm: "VLRAGGraphRLM") -> List[str]:
    """Generate sub-queries to improve recall for complex queries.
    
    Uses RLM to break down the original query into multiple
    complementary sub-queries for broader retrieval coverage.
    
    Args:
        original_query: The original user query
        rlm: Initialized RLM instance
        
    Returns:
        List of sub-queries including the original
    """
    prompt = f"""Given the following query, generate 2-3 complementary sub-queries that would help retrieve comprehensive information.

Original query: "{original_query}"

Generate sub-queries that:
1. Cover different aspects or angles of the main question
2. Use different keywords or phrasings for the same concept
3. Include both broad and specific interpretations

Respond with ONLY the sub-queries, one per line. Do not include numbering or explanations.

Example:
Query: "What are the benefits of machine learning in healthcare?"
Sub-queries:
machine learning applications in medical diagnosis
AI healthcare benefits and patient outcomes
automated diagnosis systems advantages

Your sub-queries:"""
    
    try:
        result = rlm.completion(prompt, "")
        lines = [line.strip() for line in result.response.strip().split('\n') if line.strip()]
        # Filter out any explanatory lines that might have been added
        sub_queries = [line for line in lines if len(line) > 10 and not line.lower().startswith(('sub-query', 'example', 'note', 'here are'))]
        
        # Always include the original query
        all_queries = [original_query] + sub_queries[:3]  # Max 3 sub-queries + original
        return list(dict.fromkeys(all_queries))  # Remove duplicates, preserve order
    except Exception:
        # If generation fails, return just the original query
        return [original_query]


# ── Retrieval-specific embedding instructions ──────────────────────────
# Qwen3-VL embedding quality improves significantly when the instruction
# matches the task.  "Represent the user's input." is the default for
# document ingestion; retrieval queries need a retrieval-oriented prompt.
_QUERY_INSTRUCTION = (
    "Find passages that are relevant to and answer the following query."
)
_DOCUMENT_INSTRUCTION = "Represent this document for retrieval."

# ── Knowledge-graph extraction prompt ──────────────────────────────────
_KG_EXTRACTION_PROMPT = (
    "You are a knowledge-graph extraction engine. Analyse the document below "
    "and produce a structured knowledge graph in Markdown.\n\n"
    "For every entity you find, output a bullet under **Entities** with its "
    "name, type (Person, Organisation, Concept, Technology, Location, Event, "
    "Metric, etc.), and a one-line description.\n\n"
    "For every relationship between entities, output a bullet under "
    "**Relationships** in the form:\n"
    "  - EntityA → relationship → EntityB (brief context)\n\n"
    "Group entities by type. Be exhaustive — capture every meaningful concept, "
    "term, and connection. Prefer precision over brevity."
)

# ── Accuracy-first retrieval parameters ────────────────────────────────
# These are intentionally wide — we prioritise recall and reranking
# accuracy over speed.
_DENSE_TOP_K = 50          # dense (Qwen3-VL embedding) retrieval depth
_KEYWORD_TOP_K = 50        # keyword retrieval depth
_RERANK_CANDIDATES = 30    # candidates sent to Qwen3-VL cross-attention reranker
_FINAL_RESULTS = 10        # top reranked results used as context for the RLM


def _run_vl_rag_query(
    query: str,
    *,
    store: Optional["MultimodalVectorStore"],
    reranker_vl: Optional["Qwen3VLRerankerProvider"],
    rrf: "ReciprocalRankFusion",
    fallback_reranker: "CompositeReranker",
    all_chunks: List[dict],
    knowledge_graph: str,
    context_budget: int,
    rlm: "VLRAGGraphRLM",
    fallback_hierarchy: Optional[List[str]],
    provider: str,
    resolved_model: Optional[str],
    max_depth: int,
    max_iterations: int,
    verbose: bool = True,
    rrf_weights: Optional[List[float]] = None,
    use_graph_augmented: bool = False,
    graph_hops: int = 2,
) -> Dict[str, Any]:
    """Run a single query through the full VL-RAG-Graph-RLM pipeline.

    This is the **single source of truth** for every query path — both
    ``run_analysis()`` and the interactive REPL call this function so that
    every query is guaranteed to go through:

    1. Qwen3-VL dense embedding search (with retrieval instruction)
    2. Keyword search + Reciprocal Rank Fusion
    3. Qwen3-VL cross-attention reranking (text + images)
    4. Knowledge-graph context augmentation
    5. RLM recursive completion (with provider hierarchy fallback)
    """
    sources_info: List[Dict[str, Any]] = []

    # ── Stage 1-3: Retrieval → Fusion → Reranking ──────────────────
    if store is not None and HAS_QWEN3VL:
        # 1. Dense search with retrieval-specific instruction
        dense_results = store.search(
            query, top_k=_DENSE_TOP_K, instruction=_QUERY_INSTRUCTION,
        )
        if verbose:
            print(f"    Dense search: {len(dense_results)} results")

        # 2. Keyword search + RRF fusion
        keyword_results = _keyword_search(store, query, top_k=_KEYWORD_TOP_K)
        if verbose:
            print(f"    Keyword search: {len(keyword_results)} results")

        if keyword_results:
            fused_results = rrf.fuse(
                [dense_results, keyword_results], weights=rrf_weights or [4.0, 1.0],
            )
        else:
            fused_results = dense_results
        if verbose:
            print(f"    After RRF fusion: {len(fused_results)} results")

        # 3. Qwen3-VL cross-attention reranking (text + images)
        if reranker_vl and fused_results:
            docs_for_rerank: List[Dict[str, Any]] = []
            for r in fused_results[:_RERANK_CANDIDATES]:
                doc = store.get(r.id)
                if doc:
                    doc_dict: Dict[str, Any] = {"text": doc.content}
                    if doc.image_path:
                        doc_dict["image"] = doc.image_path
                    docs_for_rerank.append(doc_dict)

            reranked_indices = reranker_vl.rerank(
                query={"text": query}, documents=docs_for_rerank,
            )
            reordered = []
            for idx, score in reranked_indices:
                if idx < len(fused_results):
                    result = fused_results[idx]
                    result.composite_score = score * 100
                    reordered.append(result)
            final_results = reordered[:_FINAL_RESULTS]
            if verbose:
                print(
                    f"    After Qwen3-VL reranking: {len(final_results)} results"
                    f" (from {len(docs_for_rerank)} candidates)"
                )
        else:
            final_results = fused_results[:_FINAL_RESULTS]

        # Assemble context from reranked results
        context_parts = []
        for i, result in enumerate(final_results):
            doc = store.get(result.id)
            if doc:
                context_parts.append(
                    f"[Source {i + 1}] (score: {result.composite_score:.1f})\n"
                    f"{doc.content}"
                )
                sources_info.append(
                    {
                        "content": doc.content[:200],
                        "score": result.composite_score,
                        "metadata": doc.metadata,
                    }
                )
        context = "\n\n---\n\n".join(context_parts)
    else:
        # Fallback: CompositeReranker without VL embeddings
        search_results = [
            SearchResult(
                id=i,
                content=chunk.get("content", ""),
                metadata={"type": chunk.get("type", "text")},
                semantic_score=1.0,
            )
            for i, chunk in enumerate(all_chunks)
        ]
        reranked, _ = fallback_reranker.process(
            query, [r.__dict__ for r in search_results],
        )
        reranked.sort(
            key=lambda x: x.get("composite_score", 0), reverse=True,
        )
        top_chunks = reranked[:_FINAL_RESULTS]
        context = "\n\n---\n\n".join(
            [c.get("content", "") for c in top_chunks]
        )

    # ── Stage 4: Knowledge-graph context augmentation ──────────────
    if knowledge_graph and not knowledge_graph.startswith("Could not"):
        kg_budget = min(8000, context_budget // 3)
        context = (
            f"[Knowledge Graph]\n{knowledge_graph[:kg_budget]}\n\n---\n\n"
            f"{context}"
        )

    # ── Stage 4b: Graph-augmented retrieval (optional) ─────────────
    if use_graph_augmented and knowledge_graph and not knowledge_graph.startswith("Could not"):
        try:
            from vl_rag_graph_rlm.graph_retrieval import augment_retrieval_with_graph
            
            # Build chunk list from context for entity extraction
            context_chunks = [{"content": context}]
            
            graph_context = augment_retrieval_with_graph(
                query=query,
                retrieved_chunks=context_chunks,
                kg_text=knowledge_graph,
                max_hops=graph_hops,
                max_expanded=15,
            )
            
            if graph_context:
                if verbose:
                    print(f"    Graph-augmented: added {graph_context.count('→')} connections")
                context = (
                    f"[Graph-Augmented Context]\n{graph_context}\n\n---\n\n"
                    f"{context}"
                )
        except Exception as e:
            if verbose:
                print(f"    Graph augmentation skipped: {e}")

    # ── Stage 5: RLM recursive completion ──────────────────────────
    try:
        q_start = time.time()
        result = rlm.completion(query, context[:context_budget])
        elapsed = time.time() - q_start
        if verbose:
            print(f"    RLM answer in {elapsed:.2f}s")
        return {
            "query": query,
            "response": result.response,
            "time": elapsed,
            "sources": sources_info,
        }
    except Exception as e:
        # Provider hierarchy fallback
        if fallback_hierarchy and len(fallback_hierarchy) > 1:
            fallback_result = _try_fallback_query(
                query, context[:context_budget], fallback_hierarchy,
                provider, resolved_model, max_depth, max_iterations,
            )
            if fallback_result:
                if verbose:
                    print(
                        f"    Fallback answer via "
                        f"{fallback_result.get('fallback_provider', '?')}"
                    )
                return fallback_result
        return {
            "query": query,
            "response": f"Error: {e}",
            "time": 0,
            "sources": [],
        }


def _resolve_provider(provider: str) -> tuple[str, bool]:
    """Resolve provider name, handling 'auto' mode.

    Returns:
        (resolved_provider_name, is_auto)
    """
    is_auto = provider == "auto"
    if is_auto:
        try:
            provider = resolve_auto_provider()
        except RuntimeError as e:
            print(f"Error: {e}")
            print("Run 'vrlmrag --show-hierarchy' to see provider status.")
            sys.exit(1)
    if provider not in SUPPORTED_PROVIDERS:
        print(f"Error: Unknown provider '{provider}'")
        print(f"Supported: {', '.join(SUPPORTED_PROVIDERS.keys())}")
        sys.exit(1)
    return provider, is_auto


def _load_knowledge_graph(kg_path: Path) -> str:
    """Load persisted knowledge graph from disk."""
    if kg_path.exists():
        return kg_path.read_text(encoding="utf-8")
    return ""


def _save_knowledge_graph(kg_path: Path, kg_text: str) -> None:
    """Persist knowledge graph to disk."""
    kg_path.parent.mkdir(parents=True, exist_ok=True)
    kg_path.write_text(kg_text, encoding="utf-8")


def _merge_knowledge_graphs(existing: str, new_fragment: str) -> str:
    """Merge a new KG fragment into the existing knowledge graph."""
    if not existing:
        return new_fragment
    if not new_fragment:
        return existing
    return f"{existing}\n\n---\n\n{new_fragment}"


def sliding_window_chunks(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks using sliding window.
    
    Args:
        text: The text to split
        chunk_size: Maximum characters per chunk
        overlap: Number of characters to overlap between chunks
        
    Returns:
        List of text chunks with overlap
    """
    if not text:
        return []
    
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = min(start + chunk_size, text_len)
        
        # Try to break at a sentence or word boundary
        if end < text_len:
            # Look for sentence ending
            for i in range(end, max(start, end - 100), -1):
                if text[i] in '.!?\n':
                    end = i + 1
                    break
            else:
                # No sentence break found, look for word boundary
                for i in range(end, max(start, end - 50), -1):
                    if text[i].isspace():
                        end = i
                        break
        
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        
        # Move window with overlap
        start = end - overlap if end < text_len else text_len
        
        # Prevent infinite loop if no progress
        if start >= end:
            break
    
    return chunks


def run_interactive_session(
    provider: str,
    input_path: Optional[str] = None,
    model: Optional[str] = None,
    max_depth: int = 3,
    max_iterations: int = 10,
    store_dir: Optional[str] = None,
    use_api: bool = False,
    text_only: bool = False,
) -> None:
    """Run an interactive VL-RAG-Graph-RLM session.

    Loads VL models once at startup and keeps them resident in memory.
    Supports incremental document addition and continuous querying.
    The knowledge graph grows across queries and persists to disk.

    Commands (type at the ``vrlmrag>`` prompt):
        <any text>            — run as a query against loaded documents
        /add <path>           — add more documents (file or folder)
        /kg                   — show the current knowledge graph
        /stats                — show session statistics
        /save [path]          — save report to file
        /help                 — show available commands
        /quit or /exit        — end session
    """
    import json as _json
    import readline  # noqa: F401 — enables arrow-key history in input()

    provider, is_auto = _resolve_provider(provider)
    prov_info = SUPPORTED_PROVIDERS[provider]
    context_budget = prov_info["context_budget"]

    api_key = os.getenv(prov_info["env_key"])
    if not api_key:
        print(f"Error: {prov_info['env_key']} not set")
        print(f"Get your API key from: {prov_info['url']}")
        sys.exit(1)

    resolved_model = model  # explicit --model flag takes priority

    # ----------------------------------------------------------------
    # Determine persistence directory
    # ----------------------------------------------------------------
    if store_dir:
        session_dir = Path(store_dir)
    elif input_path:
        p = Path(input_path)
        session_dir = (p.parent if p.is_file() else p) / ".vrlmrag_store"
    else:
        session_dir = Path.cwd() / ".vrlmrag_store"
    session_dir.mkdir(parents=True, exist_ok=True)
    storage_file = str(session_dir / "embeddings.json")
    kg_file = session_dir / "knowledge_graph.md"

    # ----------------------------------------------------------------
    # Banner
    # ----------------------------------------------------------------
    print("=" * 70)
    print(f"vrlmrag v{__version__} — Interactive Session")
    print("=" * 70)
    print(f"Provider:       {provider}" + (" (auto)" if is_auto else ""))
    print(f"Model:          {resolved_model or '(auto-detect from env/defaults)'}")
    print(f"Context budget: {context_budget:,} chars")
    print(f"Store:          {session_dir}")
    if input_path:
        print(f"Input:          {input_path}")
    print("=" * 70)

    # ----------------------------------------------------------------
    # Step 1: Load VL models (once)
    # For interactive sessions, the local model lock is held for the
    # entire session since the model stays resident in RAM.
    # API-based embeddings bypass the lock entirely.
    # ----------------------------------------------------------------
    embedding_model_name = os.getenv("VRLMRAG_LOCAL_EMBEDDING_MODEL", "Qwen/Qwen3-VL-Embedding-2B")
    reranker_model_name = os.getenv("VRLMRAG_RERANKER_MODEL", "ms-marco-MiniLM-L-12-v2")

    reranker_vl = None
    store: Optional["MultimodalVectorStore"] = None
    _session_lock_ctx = None  # held for entire interactive session

    if text_only and HAS_TEXT_EMBEDDING:
        text_model = os.getenv("VRLMRAG_TEXT_ONLY_MODEL", "Qwen/Qwen3-Embedding-0.6B")
        _session_lock_ctx = local_model_lock(text_model, description="CLI interactive (text-only)")
        _session_lock_ctx.__enter__()
        print(f"\n[init] Loading text-only embedding ({text_model})...")
        embedder = create_text_embedder(model_name=text_model)
        embedding_model_name = text_model
        print(f"  Embedding dim: {embedder.embedding_dim}")

        storage_file_text = str(Path(session_dir) / "embeddings_text.json")
        store = MultimodalVectorStore(
            embedding_provider=embedder,
            storage_path=storage_file_text,
        )
        existing_docs = len(store.documents)
        if existing_docs:
            print(f"  Loaded {existing_docs} existing documents from store")

        if HAS_FLASHRANK:
            print("[init] Loading FlashRank Reranker...")
            reranker_vl = create_flashrank_reranker(model_name=reranker_model_name)
            print(f"  Reranker loaded ({reranker_model_name})")
    elif use_api and HAS_API_EMBEDDING:
        print("\n[init] Loading API-based embedding (OpenRouter + ZenMux)...")
        embedder = create_api_embedder()
        embedding_model_name = os.getenv("VRLMRAG_EMBEDDING_MODEL", "openai/text-embedding-3-small")
        print(f"  Embedding model: {embedding_model_name}")

        store = MultimodalVectorStore(
            embedding_provider=embedder,
            storage_path=storage_file,
        )
        existing_docs = len(store.documents)
        if existing_docs:
            print(f"  Loaded {existing_docs} existing documents from store")

        if HAS_FLASHRANK:
            print("[init] Loading FlashRank Reranker...")
            reranker_vl = create_flashrank_reranker(model_name=reranker_model_name)
            print(f"  Reranker loaded ({reranker_model_name})")
    elif HAS_QWEN3VL:
        _session_lock_ctx = local_model_lock(embedding_model_name, description="CLI interactive (Qwen3-VL)")
        _session_lock_ctx.__enter__()

        import torch

        device = "mps" if torch.backends.mps.is_available() else "cpu"

        # Load embedder → create store (stays loaded for /add commands)
        print(f"\n[init] Loading Qwen3-VL Embedding model ({device})...")
        embedder = create_qwen3vl_embedder(
            model_name=embedding_model_name, device=device
        )
        print(f"  Embedding dim: {embedder.embedding_dim}")

        store = MultimodalVectorStore(
            embedding_provider=embedder,
            storage_path=storage_file,
        )
        existing_docs = len(store.documents)
        if existing_docs:
            print(f"  Loaded {existing_docs} existing documents from store")

        if HAS_FLASHRANK:
            print("[init] Loading FlashRank Reranker...")
            reranker_vl = create_flashrank_reranker(model_name=reranker_model_name)
            print(f"  Reranker loaded ({reranker_model_name})")
        else:
            print("[init] FlashRank not available — using RRF-only retrieval")
            print("  Install with: pip install flashrank")
    else:
        print("\n[init] No embedding provider available — using fallback text reranking")
        print("  Install with: pip install torch transformers qwen-vl-utils torchvision")

    # ----------------------------------------------------------------
    # Step 2: Initialize RLM
    # ----------------------------------------------------------------
    rlm_kwargs: Dict[str, Any] = {
        "provider": provider,
        "temperature": 0.0,
        "max_depth": max_depth,
        "max_iterations": max_iterations,
    }
    if resolved_model:
        rlm_kwargs["model"] = resolved_model

    print(f"[init] Initializing RLM ({provider})...")
    rlm = VLRAGGraphRLM(**rlm_kwargs)
    print(f"  Model: {rlm.model}")

    fallback_hierarchy = get_available_providers() if is_auto else None

    # ----------------------------------------------------------------
    # Step 3: Load persisted knowledge graph
    # ----------------------------------------------------------------
    knowledge_graph = _load_knowledge_graph(kg_file)
    if knowledge_graph:
        print(f"[init] Loaded existing knowledge graph ({len(knowledge_graph):,} chars)")

    # ----------------------------------------------------------------
    # Step 4: Process initial documents (if provided)
    # Local mode: Parakeet ASR transcribes audio → text chunks
    # API mode: audio/video handled by omni model at embedding time
    # ----------------------------------------------------------------
    _transcriber = None
    if not use_api and not text_only and HAS_PARAKEET:
        _transcriber = create_parakeet_transcriber(
            cache_dir=str(Path.home() / ".vrlmrag" / "parakeet_cache"),
        )
    processor = DocumentProcessor(
        transcription_provider=_transcriber,
        use_api=use_api,
    )
    all_chunks: List[dict] = []
    all_documents: List[dict] = []
    rrf = ReciprocalRankFusion(k=60)
    fallback_reranker = CompositeReranker()
    query_count = 0
    total_query_time = 0.0

    def _ingest_path(path_str: str) -> int:
        """Ingest documents from a path, returning count of new embeddings.

        The embedder stays loaded permanently and the lightweight
        FlashRank reranker (~34 MB) coexists with it — no model
        swapping needed.
        """
        nonlocal knowledge_graph
        docs = processor.process_path(path_str)
        if isinstance(docs, dict):
            docs = [docs]

        new_chunks: List[dict] = []
        for doc in docs:
            new_chunks.extend(doc.get("chunks", []))
        all_documents.extend(docs)
        all_chunks.extend(new_chunks)

        embedded = 0
        if store is not None:
            for i, chunk in enumerate(new_chunks):
                content = chunk.get("content", "")
                if not content.strip():
                    continue
                metadata = {"type": chunk.get("type", "text")}
                if "slide" in chunk:
                    metadata["slide"] = chunk["slide"]
                store.add_text(
                    content=content, metadata=metadata,
                    instruction=_DOCUMENT_INSTRUCTION,
                )
                embedded += 1

            # Embed images
            for doc in docs:
                for img_info in doc.get("image_data", []):
                    try:
                        temp_path = f"/tmp/vrlmrag_{img_info['filename']}"
                        with open(temp_path, "wb") as f:
                            f.write(img_info["blob"])
                        store.add_image(
                            image_path=temp_path,
                            description=f"Image from slide {img_info['slide']}",
                            metadata={
                                "type": "image",
                                "slide": img_info["slide"],
                                "filename": img_info["filename"],
                            },
                            instruction=_DOCUMENT_INSTRUCTION,
                        )
                        embedded += 1
                    except Exception as e:
                        print(f"  Warning: Could not embed image {img_info['filename']}: {e}")

            # Embed video frames from media documents
            for doc in docs:
                frame_paths = doc.get("frame_paths", [])
                if frame_paths:
                    print(f"  Embedding {len(frame_paths)} video frames...")
                    for i, frame_path in enumerate(frame_paths):
                        try:
                            store.add_image(
                                image_path=frame_path,
                                description=f"Video frame {i+1} from {Path(doc['path']).name}",
                                metadata={
                                    "type": "video_frame",
                                    "frame_index": i,
                                    "source": doc["path"],
                                },
                                instruction=_DOCUMENT_INSTRUCTION,
                            )
                            embedded += 1
                        except Exception as e:
                            print(f"  Warning: Could not embed frame {i}: {e}")

        # Build/extend knowledge graph for new documents
        if new_chunks:
            kg_char_limit = min(context_budget, 25000)
            kg_doc_limit = max(2000, kg_char_limit // max(len(docs), 1))
            kg_context = "\n\n".join(
                [d.get("content", "")[:kg_doc_limit] for d in docs]
            )
            try:
                kg_result = rlm.completion(
                    _KG_EXTRACTION_PROMPT,
                    kg_context[:kg_char_limit],
                )
                knowledge_graph = _merge_knowledge_graphs(
                    knowledge_graph, kg_result.response
                )
                _save_knowledge_graph(kg_file, knowledge_graph)
            except Exception as e:
                print(f"  Warning: Could not extend knowledge graph: {e}")

        print(f"  Processed {len(docs)} doc(s), {len(new_chunks)} chunks, {embedded} embedded")
        return embedded

    if input_path:
        print(f"\n[ingest] Processing: {input_path}")
        _ingest_path(input_path)

    # ----------------------------------------------------------------
    # Step 5: Query helper (delegates to shared VL-RAG pipeline)
    # ----------------------------------------------------------------
    def _run_query(q: str) -> Dict[str, Any]:
        """Run a single query through the full VL-RAG-Graph-RLM pipeline."""
        nonlocal query_count, total_query_time
        result = _run_vl_rag_query(
            q,
            store=store,
            reranker_vl=reranker_vl,
            rrf=rrf,
            fallback_reranker=fallback_reranker,
            all_chunks=all_chunks,
            knowledge_graph=knowledge_graph,
            context_budget=context_budget,
            rlm=rlm,
            fallback_hierarchy=fallback_hierarchy,
            provider=provider,
            resolved_model=resolved_model,
            max_depth=max_depth,
            max_iterations=max_iterations,
            verbose=False,
        )
        query_count += 1
        total_query_time += result.get("time", 0)
        return result

    # ----------------------------------------------------------------
    # Step 6: Interactive REPL
    # The try/finally ensures the session lock is released on exit.
    # ----------------------------------------------------------------
    print("\n" + "=" * 70)
    print("Interactive session ready. Type queries or commands.")
    print("Commands: /add <path>  /kg  /stats  /save [path]  /help  /quit")
    print("=" * 70)

    try:
        while True:
            try:
                user_input = input("\nvrlmrag> ").strip()
            except (EOFError, KeyboardInterrupt):
                print("\n[session] Goodbye.")
                break

            if not user_input:
                continue

            # ---- Commands ----
            if user_input.lower() in ("/quit", "/exit", "/q"):
                print("[session] Goodbye.")
                break

            if user_input.lower() == "/help":
                print(
                    "Commands:\n"
                    "  <text>            Query loaded documents\n"
                    "  /add <path>       Add more documents (file or folder)\n"
                    "  /kg               Show the current knowledge graph\n"
                    "  /stats            Show session statistics\n"
                    "  /save [path]      Save current KG + session report\n"
                    "  /help             Show this help\n"
                    "  /quit             End session"
                )
                continue

            if user_input.lower() == "/kg":
                if knowledge_graph:
                    print(f"\n--- Knowledge Graph ({len(knowledge_graph):,} chars) ---")
                    print(knowledge_graph[:8000])
                    if len(knowledge_graph) > 8000:
                        print(f"\n... ({len(knowledge_graph) - 8000:,} more chars)")
                else:
                    print("No knowledge graph built yet. Add documents first.")
                continue

            if user_input.lower() == "/stats":
                doc_count = len(store.documents) if store else len(all_chunks)
                print(
                    f"\n--- Session Statistics ---\n"
                    f"  Provider:           {provider} (model: {rlm.model})\n"
                    f"  Documents in store: {doc_count}\n"
                    f"  Total chunks:       {len(all_chunks)}\n"
                    f"  Knowledge graph:    {len(knowledge_graph):,} chars\n"
                    f"  Queries run:        {query_count}\n"
                    f"  Total query time:   {total_query_time:.2f}s\n"
                    f"  Avg query time:     {(total_query_time / query_count):.2f}s"
                    if query_count else
                    f"\n--- Session Statistics ---\n"
                    f"  Provider:           {provider} (model: {rlm.model})\n"
                    f"  Documents in store: {doc_count}\n"
                    f"  Total chunks:       {len(all_chunks)}\n"
                    f"  Knowledge graph:    {len(knowledge_graph):,} chars\n"
                    f"  Queries run:        0\n"
                    f"  Total query time:   0.00s"
                )
                continue

            if user_input.lower().startswith("/save"):
                parts = user_input.split(maxsplit=1)
                save_path = parts[1] if len(parts) > 1 else str(session_dir / "session_report.md")
                generator = MarkdownReportGenerator()
                results = {
                    "provider": provider,
                    "model": rlm.model,
                    "embedding_model": embedding_model_name if HAS_QWEN3VL else "N/A",
                    "reranker_model": reranker_model_name if HAS_QWEN3VL else "N/A",
                    "document_count": len(all_documents),
                    "total_chunks": len(all_chunks),
                    "embedded_count": len(store.documents) if store else 0,
                    "execution_time": total_query_time,
                    "documents": all_documents,
                    "knowledge_graph": knowledge_graph,
                    "queries": [],
                }
                report = generator.generate(results, save_path)
                print(f"  Report saved to: {save_path}")
                continue

            if user_input.lower().startswith("/add"):
                parts = user_input.split(maxsplit=1)
                if len(parts) < 2:
                    print("Usage: /add <path>")
                    continue
                add_path = parts[1].strip()
                if not Path(add_path).exists():
                    print(f"Error: Path not found: {add_path}")
                    continue
                print(f"[ingest] Processing: {add_path}")
                _ingest_path(add_path)
                continue

            # ---- Default: treat as query ----
            print(f"  Querying...")
            result = _run_query(user_input)
            print(f"\n{result['response']}")
            if result.get("time"):
                print(f"\n  [{result['time']:.2f}s]")
            if result.get("sources"):
                print(f"  [{len(result['sources'])} sources retrieved]")
    finally:
        # Release the local model lock when the session ends
        if _session_lock_ctx is not None:
            _session_lock_ctx.__exit__(None, None, None)


# ── Collection operations ──────────────────────────────────────────────


def run_collection_add(
    collection_names: List[str],
    input_path: str,
    provider: str,
    model: Optional[str] = None,
    max_depth: int = 3,
    max_iterations: int = 10,
    description: str = "",
    use_api: bool = False,
    text_only: bool = False,
) -> None:
    """Add documents from *input_path* into one or more named collections.

    Creates the collection if it does not exist.  Embeds new content with
    Qwen3-VL, builds/merges the knowledge graph, and persists everything.
    """
    provider, is_auto = _resolve_provider(provider)
    prov_info = SUPPORTED_PROVIDERS[provider]
    context_budget = prov_info["context_budget"]

    api_key = os.getenv(prov_info["env_key"])
    if not api_key:
        print(f"Error: {prov_info['env_key']} not set")
        sys.exit(1)

    resolved_model = model

    # Process documents once
    _transcriber = None
    if not use_api and not text_only and HAS_PARAKEET:
        _transcriber = create_parakeet_transcriber(
            cache_dir=str(Path.home() / ".vrlmrag" / "parakeet_cache"),
        )
    processor = DocumentProcessor(
        transcription_provider=_transcriber,
        use_api=use_api,
    )
    print(f"[collection] Processing: {input_path}")
    documents = processor.process_path(input_path)
    if isinstance(documents, dict):
        documents = [documents]

    all_chunks: List[dict] = []
    for doc in documents:
        all_chunks.extend(doc.get("chunks", []))
    print(f"  {len(documents)} document(s), {len(all_chunks)} chunks")

    # Initialise RLM for KG extraction
    rlm_kwargs: Dict[str, Any] = {
        "provider": provider,
        "temperature": 0.0,
        "max_depth": max_depth,
        "max_iterations": max_iterations,
    }
    if resolved_model:
        rlm_kwargs["model"] = resolved_model
    rlm = VLRAGGraphRLM(**rlm_kwargs)

    for cname in collection_names:
        meta = create_collection(cname, description=description)
        slug = meta["name"]
        storage_file = collection_embeddings_path(slug)
        kg_file = collection_kg_path(slug)

        print(f"\n[collection:{slug}] Adding documents...")

        embedder = None
        store = None
        embedded_count = 0
        skipped_count = 0
        _coll_lock_ctx = None

        if text_only and HAS_TEXT_EMBEDDING:
            text_model = os.getenv("VRLMRAG_TEXT_ONLY_MODEL", "Qwen/Qwen3-Embedding-0.6B")
            _coll_lock_ctx = local_model_lock(text_model, description=f"CLI collection_add:{slug} (text-only)")
            _coll_lock_ctx.__enter__()
            embedder = create_text_embedder(model_name=text_model)
            store = MultimodalVectorStore(
                embedding_provider=embedder, storage_path=storage_file,
            )
        elif use_api and HAS_API_EMBEDDING:
            embedder = create_api_embedder()
            store = MultimodalVectorStore(
                embedding_provider=embedder, storage_path=storage_file,
            )
        elif HAS_QWEN3VL:
            _emb_model = os.getenv("VRLMRAG_LOCAL_EMBEDDING_MODEL", "Qwen/Qwen3-VL-Embedding-2B")
            _coll_lock_ctx = local_model_lock(_emb_model, description=f"CLI collection_add:{slug} (Qwen3-VL)")
            _coll_lock_ctx.__enter__()

            import torch

            device = "mps" if torch.backends.mps.is_available() else "cpu"
            embedder = create_qwen3vl_embedder(
                model_name=_emb_model,
                device=device,
            )
            store = MultimodalVectorStore(
                embedding_provider=embedder, storage_path=storage_file,
            )
            existing = len(store.documents)
            if existing:
                print(f"  Loaded {existing} existing embeddings")

            for chunk in all_chunks:
                content = chunk.get("content", "")
                if not content.strip():
                    continue
                if store.content_exists(content):
                    skipped_count += 1
                    continue
                metadata = {"type": chunk.get("type", "text")}
                if "slide" in chunk:
                    metadata["slide"] = chunk["slide"]
                store.add_text(
                    content=content, metadata=metadata,
                    instruction=_DOCUMENT_INSTRUCTION,
                )
                embedded_count += 1

            # Embed images
            for doc in documents:
                for img_info in doc.get("image_data", []):
                    try:
                        temp_path = f"/tmp/vrlmrag_{img_info['filename']}"
                        with open(temp_path, "wb") as f:
                            f.write(img_info["blob"])
                        prev = len(store.documents)
                        store.add_image(
                            image_path=temp_path,
                            description=f"Image from slide {img_info['slide']}",
                            metadata={
                                "type": "image",
                                "slide": img_info["slide"],
                                "filename": img_info["filename"],
                            },
                            instruction=_DOCUMENT_INSTRUCTION,
                        )
                        if len(store.documents) > prev:
                            embedded_count += 1
                        else:
                            skipped_count += 1
                    except Exception as e:
                        print(f"  Warning: image embed failed: {e}")

            print(f"  New: {embedded_count} | Skipped: {skipped_count} | Total: {len(store.documents)}")
        else:
            print("  Warning: Qwen3-VL not available — skipping embeddings")

        # Release the local model lock — KG extraction below is API-based
        if _coll_lock_ctx is not None:
            _coll_lock_ctx.__exit__(None, None, None)
            _coll_lock_ctx = None

        # Build/merge knowledge graph
        knowledge_graph = collection_load_kg(slug)
        if all_chunks:
            kg_char_limit = min(context_budget, 25000)
            kg_doc_limit = max(2000, kg_char_limit // max(len(documents), 1))
            kg_context = "\n\n".join(
                [d.get("content", "")[:kg_doc_limit] for d in documents]
            )
            try:
                kg_result = rlm.completion(_KG_EXTRACTION_PROMPT, kg_context[:kg_char_limit])
                knowledge_graph = collection_merge_kg(knowledge_graph, kg_result.response)
                collection_save_kg(slug, knowledge_graph)
                print(f"  Knowledge graph: {len(knowledge_graph):,} chars")
            except Exception as e:
                print(f"  Warning: KG extraction failed: {e}")

        # Determine embedding model name for tracking
        embedding_model_name = ""
        if text_only and HAS_TEXT_EMBEDDING:
            embedding_model_name = os.getenv("VRLMRAG_TEXT_ONLY_MODEL", "Qwen/Qwen3-Embedding-0.6B")
        elif use_api and HAS_API_EMBEDDING:
            embedding_model_name = os.getenv("VRLMRAG_EMBEDDING_MODEL", "openai/text-embedding-3-small")
        elif HAS_QWEN3VL:
            embedding_model_name = os.getenv("VRLMRAG_LOCAL_EMBEDDING_MODEL", "Qwen/Qwen3-VL-Embedding-2B")
        
        reranker_model_name = os.getenv("VRLMRAG_RERANKER_MODEL", "ms-marco-MiniLM-L-12-v2")
        
        record_source(
            slug, input_path, len(documents), len(all_chunks),
            embedding_model=embedding_model_name,
            reranker_model=reranker_model_name,
        )
        print(f"  Collection '{slug}' updated.")


def run_collection_query(
    collection_names: List[str],
    query: str,
    provider: str,
    model: Optional[str] = None,
    max_depth: int = 3,
    max_iterations: int = 10,
    output: Optional[str] = None,
    use_api: bool = False,
    text_only: bool = False,
) -> None:
    """Query one or more collections, blending their stores and KGs.

    This is fully scriptable — no user interaction required.
    """
    provider, is_auto = _resolve_provider(provider)
    prov_info = SUPPORTED_PROVIDERS[provider]
    context_budget = prov_info["context_budget"]

    api_key = os.getenv(prov_info["env_key"])
    if not api_key:
        print(f"Error: {prov_info['env_key']} not set")
        sys.exit(1)

    resolved_model = model
    fallback_hierarchy = get_available_providers() if is_auto else None

    # Initialise RLM
    rlm_kwargs: Dict[str, Any] = {
        "provider": provider,
        "temperature": 0.0,
        "max_depth": max_depth,
        "max_iterations": max_iterations,
    }
    if resolved_model:
        rlm_kwargs["model"] = resolved_model
    rlm = VLRAGGraphRLM(**rlm_kwargs)

    # Load and blend collections
    # Local models acquire the cross-process lock to ensure only one
    # model is loaded in RAM at a time (across all CLI/MCP sessions).
    blended_store = None
    blended_kg = ""
    all_chunks: List[dict] = []
    embedder = None
    reranker_vl = None
    _cq_lock_ctx = None

    if text_only and HAS_TEXT_EMBEDDING:
        text_model = os.getenv("VRLMRAG_TEXT_ONLY_MODEL", "Qwen/Qwen3-Embedding-0.6B")
        _cq_lock_ctx = local_model_lock(text_model, description="CLI collection_query (text-only)")
        _cq_lock_ctx.__enter__()
        embedder = create_text_embedder(model_name=text_model)
    elif use_api and HAS_API_EMBEDDING:
        embedder = create_api_embedder()
    elif HAS_QWEN3VL:
        _emb_model = os.getenv("VRLMRAG_LOCAL_EMBEDDING_MODEL", "Qwen/Qwen3-VL-Embedding-2B")
        _cq_lock_ctx = local_model_lock(_emb_model, description="CLI collection_query (Qwen3-VL)")
        _cq_lock_ctx.__enter__()

        import torch

        device = "mps" if torch.backends.mps.is_available() else "cpu"
        embedder = create_qwen3vl_embedder(
            model_name=_emb_model,
            device=device,
        )

    collection_labels = []
    for cname in collection_names:
        if not collection_exists(cname):
            print(f"Error: Collection '{cname}' does not exist")
            print("Run 'vrlmrag --collection-list' to see available collections")
            sys.exit(1)

        meta = load_collection_meta(cname)
        slug = meta["name"]
        collection_labels.append(slug)

        # Merge KG
        kg = collection_load_kg(slug)
        if kg:
            blended_kg = collection_merge_kg(blended_kg, kg)

        # Load store
        if embedder:
            storage_file = collection_embeddings_path(slug)
            store = MultimodalVectorStore(
                embedding_provider=embedder, storage_path=storage_file,
            )
            if blended_store is None:
                blended_store = store
            else:
                for doc_id, doc in store.documents.items():
                    if doc_id not in blended_store.documents:
                        blended_store.documents[doc_id] = doc

    # Load lightweight FlashRank reranker (~34 MB — coexists with embedder)
    if HAS_FLASHRANK:
        reranker_vl = create_flashrank_reranker(
            model_name=os.getenv("VRLMRAG_RERANKER_MODEL", "ms-marco-MiniLM-L-12-v2"),
        )

    # Release the local model lock — RLM query below is API-based
    if _cq_lock_ctx is not None:
        _cq_lock_ctx.__exit__(None, None, None)
        _cq_lock_ctx = None

    label = " + ".join(collection_labels)
    print(f"[query] Collections: {label}")
    if blended_store:
        print(f"  Total documents in blended store: {len(blended_store.documents)}")
    if blended_kg:
        print(f"  Blended KG: {len(blended_kg):,} chars")

    rrf = ReciprocalRankFusion(k=60)
    fallback_reranker = CompositeReranker()

    print(f"\n  Query: {query}")
    result = _run_vl_rag_query(
        query,
        store=blended_store,
        reranker_vl=reranker_vl,
        rrf=rrf,
        fallback_reranker=fallback_reranker,
        all_chunks=all_chunks,
        knowledge_graph=blended_kg,
        context_budget=context_budget,
        rlm=rlm,
        fallback_hierarchy=fallback_hierarchy,
        provider=provider,
        resolved_model=resolved_model,
        max_depth=max_depth,
        max_iterations=max_iterations,
    )

    print(f"\n{result['response']}")
    if result.get("time"):
        print(f"\n  [{result['time']:.2f}s]")
    if result.get("sources"):
        print(f"  [{len(result['sources'])} sources retrieved]")

    # Optionally save to file
    if output:
        report_lines = [
            f"# Collection Query: {label}",
            "",
            f"**Query:** {query}",
            f"**Provider:** {provider}",
            f"**Model:** {rlm.model}",
            f"**Time:** {result.get('time', 0):.2f}s",
            "",
            "## Response",
            "",
            result["response"],
            "",
        ]
        if result.get("sources"):
            report_lines.append("## Sources")
            report_lines.append("")
            for src in result["sources"]:
                score = src.get("score", 0)
                preview = src.get("content", "")[:120]
                report_lines.append(f"- [Score: {score:.2f}] {preview}...")
            report_lines.append("")

        Path(output).write_text("\n".join(report_lines), encoding="utf-8")
        print(f"\n  Report saved to: {output}")


def show_collection_list() -> None:
    """Print all available collections."""
    collections = _list_collections_meta()
    if not collections:
        print("No collections found.")
        print(f"Create one with: vrlmrag -c <name> --add <path>")
        return

    print(f"{'Name':<25} {'Docs':>6} {'Chunks':>8} {'Sources':>8}  Updated")
    print("-" * 80)
    for meta in collections:
        name = meta.get("display_name", meta["name"])
        docs = meta.get("document_count", 0)
        chunks = meta.get("chunk_count", 0)
        sources = len(meta.get("sources", []))
        updated = meta.get("updated", "?")[:19]
        print(f"{name:<25} {docs:>6} {chunks:>8} {sources:>8}  {updated}")


def show_collection_info(name: str) -> None:
    """Print detailed info for a single collection."""
    if not collection_exists(name):
        print(f"Error: Collection '{name}' does not exist")
        sys.exit(1)

    meta = load_collection_meta(name)
    slug = meta["name"]

    print(f"Collection: {meta.get('display_name', slug)}")
    print(f"  Slug:        {slug}")
    print(f"  Description: {meta.get('description', '(none)')}")
    print(f"  Created:     {meta.get('created', '?')[:19]}")
    print(f"  Updated:     {meta.get('updated', '?')[:19]}")
    print(f"  Documents:   {meta.get('document_count', 0)}")
    print(f"  Chunks:      {meta.get('chunk_count', 0)}")
    print(f"  Directory:   {_collection_dir(slug)}")
    
    # Embedding model info
    emb_model = meta.get('embedding_model', '')
    reranker_model = meta.get('reranker_model', '')
    print(f"  Embedding:   {emb_model or '(not tracked)'}")
    if reranker_model:
        print(f"  Reranker:    {reranker_model}")
    
    # Model history
    model_history = meta.get('model_history', [])
    if model_history:
        print(f"\n  Model change history:")
        for change in model_history[-5:]:  # Show last 5 changes
            print(f"    {change['timestamp'][:19]}: {change['previous_model']} → {change['new_model']}")
    
    # Check for mixed models in sources
    sources = meta.get('sources', [])
    source_models = set(s.get('embedding_model', '') for s in sources if s.get('embedding_model'))
    if len(source_models) > 1:
        print(f"\n  [Warning] Collection has mixed embedding models:")
        for model in sorted(source_models):
            print(f"    - {model}")
        print(f"    Run 'vrlmrag -c {slug} --reindex' to normalize with current model")

    # Embedding count
    emb_file = Path(collection_embeddings_path(slug))
    if emb_file.exists():
        import json as _json
        try:
            data = _json.loads(emb_file.read_text(encoding="utf-8"))
            print(f"  Embeddings:  {len(data.get('documents', {}))}")
        except Exception:
            print(f"  Embeddings:  (file exists, could not parse)")
    else:
        print(f"  Embeddings:  0")

    # KG size
    kg = collection_load_kg(slug)
    print(f"  KG size:     {len(kg):,} chars" if kg else "  KG size:     0 chars")

    # Sources
    sources = meta.get("sources", [])
    if sources:
        print(f"\n  Sources ({len(sources)}):")
        for src in sources:
            print(f"    - {src['path']}  ({src['documents']} docs, {src['chunks']} chunks, {src['added'][:19]})")


def main():
    provider_names = ", ".join(SUPPORTED_PROVIDERS.keys())

    parser = argparse.ArgumentParser(
        prog="vrlmrag",
        description=(
            "vrlmrag — Full VL-RAG-Graph-RLM document analysis pipeline.\n\n"
            "Process documents (PPTX, PDF, TXT, MD, MP4, WAV, MP3) through the\n"
            "complete 6-pillar multimodal pipeline: VL embeddings → RAG → reranking\n"
            "→ knowledge graph → recursive LLM reasoning → markdown report.\n\n"
            "Default: API mode (uses provider hierarchy for embeddings + LLM).\n"
            "Use --local to opt into local Qwen3-VL models (blocked for video/audio)."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "examples:\n"
            "  vrlmrag presentation.pptx                           # API mode (default)\n"
            "  vrlmrag --local presentation.pptx                   # local Qwen3-VL models\n"
            "  vrlmrag --offline presentation.pptx               # offline mode (no APIs)\n"
            "  vrlmrag --provider sambanova presentation.pptx      # explicit provider\n"
            "  vrlmrag --provider nebius document.pdf -o report.md  # with output\n"
            "  vrlmrag ./docs -q 'Summarize key findings'          # auto + query\n"
            "  vrlmrag video.mp4 -q 'What is discussed?'           # video (API only)\n"
            "  vrlmrag --interactive presentation.pptx             # interactive session\n"
            "  vrlmrag -i ./codebase                               # load & query continuously\n"
            "  vrlmrag --show-hierarchy                            # see fallback order\n"
            "  vrlmrag --list-providers                            # see all providers\n"
            "\n"
            "collections (named persistent knowledge stores):\n"
            "  vrlmrag -c research --add ./papers/              # add docs to collection\n"
            "  vrlmrag -c research -q 'Key findings?'           # query a collection\n"
            "  vrlmrag -c research -c code -q 'How implemented?' # blend collections\n"
            "  vrlmrag -c research -i                           # interactive w/ collection\n"
            "  vrlmrag --collection-list                        # list all collections\n"
            "  vrlmrag -c research --collection-info            # show collection details\n"
            "  vrlmrag -c research --collection-delete          # delete a collection\n"
            "\n"
            "model upgrade / refresh workflows:\n"
            "  vrlmrag ./docs --reindex                          # re-embed all docs with new model\n"
            "  vrlmrag ./docs --rebuild-kg                       # regenerate KG with current RLM\n"
            "  vrlmrag ./docs --reindex --rebuild-kg             # full refresh: embeddings + KG\n"
            "  vrlmrag -c research --reindex                     # reindex a collection\n"
            "\n"
            "interactive mode:\n"
            "  Loads VL models once, then lets you query continuously and add more\n"
            "  documents without reloading. Knowledge graph persists across queries.\n"
            "  Commands: /add <path>  /kg  /stats  /save  /help  /quit\n"
            "\n"
            "offline mode:\n"
            "  Use --offline when no API providers are available. Uses local Qwen3-VL\n"
            "  embeddings. Video/audio processing REQUIRES API mode and will fail.\n"
            "\n"
            "provider hierarchy (auto mode):\n"
            "  When --provider is omitted, providers are tried in PROVIDER_HIERARCHY order.\n"
            "  Default: sambanova → nebius → groq → cerebras → zai → zenmux → ...\n"
            "  Customize in .env: PROVIDER_HIERARCHY=groq,cerebras,openrouter,...\n"
            "\n"
            "video/audio safety:\n"
            "  Video/audio files FORCE API mode — local vision models are blocked.\n"
            "  Attempting video with --local or --offline will fail with a helpful error.\n"
            "\n"
            "backward-compatible aliases:\n"
            "  vrlmrag --samba-nova presentation.pptx\n"
            "  vrlmrag --nebius document.pdf\n"
            f"\nsupported providers: auto, {provider_names}\n"
            "\nSet API keys in .env or via environment variables. See .env.example."
        ),
    )

    parser.add_argument(
        "--version", "-V", action="version", version=f"%(prog)s {__version__}"
    )

    parser.add_argument(
        "--list-providers", action="store_true",
        help="List all supported providers and their API key status",
    )

    parser.add_argument(
        "--show-hierarchy", action="store_true",
        help="Show the provider fallback hierarchy and availability",
    )

    parser.add_argument(
        "--lock-status", action="store_true",
        help="Show the status of the cross-process local model lock",
    )

    parser.add_argument(
        "--provider", default="auto",
        help=f"LLM provider to use (default: auto — uses hierarchy). Options: auto, {provider_names}",
    )

    parser.add_argument(
        "input", nargs="?", metavar="PATH",
        help="File or folder to process (PPTX, PDF, TXT, MD, MP4, WAV, MP3)",
    )

    parser.add_argument(
        "--query", "-q",
        help="Custom query (default: auto-generated summary queries)",
    )

    parser.add_argument(
        "--output", "-o",
        help="Output markdown report path (default: print to stdout)",
    )

    parser.add_argument(
        "--format",
        choices=["markdown", "json"], default="markdown",
        help="Output format: markdown (default) or json (machine-readable)",
    )

    parser.add_argument(
        "--verbose", "-v", action="store_true",
        help="Enable verbose output (detailed progress information)",
    )

    parser.add_argument(
        "--quiet", "-q", action="store_true",
        help="Suppress all non-error output (silent mode)",
    )

    parser.add_argument(
        "--profile",
        choices=["fast", "balanced", "thorough", "comprehensive"],
        default="comprehensive",
        help="Analysis depth — comprehensive (default) for full VL-RAG-Graph-RLM, or fast for quick search",
    )

    parser.add_argument(
        "--comprehensive", action="store_true",
        help="Explicitly enable comprehensive mode (redundant, this is the default)",
    )

    parser.add_argument(
        "--model", "-m",
        help="Override the default model for the chosen provider",
    )

    parser.add_argument(
        "--max-depth", type=int, default=3,
        help="Maximum RLM recursion depth (default: 3)",
    )

    parser.add_argument(
        "--max-iterations", type=int, default=10,
        help="Maximum RLM iterations per call (default: 10)",
    )

    parser.add_argument(
        "--interactive", "-i", action="store_true",
        help="Start an interactive session (load VL models once, query continuously)",
    )

    parser.add_argument(
        "--store-dir",
        help="Directory for persisting embeddings and knowledge graph (default: .vrlmrag_store next to input)",
    )

    parser.add_argument(
        "--local", action="store_true",
        default=os.environ.get("VRLMRAG_LOCAL", "").lower() in ("true", "1", "yes"),
        help="Use local Qwen3-VL model instead of API embeddings (env: VRLMRAG_LOCAL). "
             "Default is API mode. Local mode is blocked for video/audio files.",
    )

    parser.add_argument(
        "--offline", action="store_true",
        help="Force offline mode — uses local Qwen3-VL embeddings when API providers unavailable. "
             "WARNING: Video/audio processing requires API mode and will fail in offline mode.",
    )

    parser.add_argument(
        "--text-only", action="store_true",
        default=os.environ.get("VRLMRAG_TEXT_ONLY", "").lower() in ("true", "1", "yes"),
        help="Use lightweight text-only embeddings — skips image/video (env: VRLMRAG_TEXT_ONLY)",
    )

    # ── Chunking configuration ────────────────────────────────────────
    parser.add_argument(
        "--chunk-size", type=int, default=1000,
        help="Maximum characters per chunk for sliding window chunking (default: 1000)",
    )

    parser.add_argument(
        "--chunk-overlap", type=int, default=200,
        help="Character overlap between chunks for sliding window (default: 200)",
    )

    parser.add_argument(
        "--use-sqlite", action="store_true",
        help="Use SQLite backend for vector store (default: JSON files)",
    )

    # ── RRF Configuration ─────────────────────────────────────────────
    parser.add_argument(
        "--rrf-dense-weight", type=float, default=4.0,
        help="Weight for dense (embedding) results in RRF fusion (default: 4.0)",
    )

    parser.add_argument(
        "--rrf-keyword-weight", type=float, default=1.0,
        help="Weight for keyword (BM25) results in RRF fusion (default: 1.0)",
    )

    parser.add_argument(
        "--multi-query", action="store_true",
        help="Multi-query retrieval — generate sub-queries to improve recall",
    )

    # ── Model upgrade / refresh flags ─────────────────────────────────
    parser.add_argument(
        "--reindex", action="store_true",
        help="Force re-embedding of all documents — bypasses cache, rebuilds vector store with current model",
    )

    parser.add_argument(
        "--rebuild-kg", action="store_true",
        help="Regenerate knowledge graph with current RLM model — re-extracts entities/relationships",
    )

    parser.add_argument(
        "--model-compare", metavar="OLD_MODEL",
        help="Compare embeddings: run query with both OLD_MODEL and current model, show divergence",
    )

    parser.add_argument(
        "--check-model", metavar="MODEL_NAME",
        help="Check if collection is compatible with target embedding model (requires -c)",
    )

    parser.add_argument(
        "--quality-check", action="store_true",
        help="RLM-powered embedding quality assessment — evaluate retrieval quality for a collection (requires -c)",
    )

    parser.add_argument(
        "--export-graph", metavar="PATH",
        help="Export knowledge graph to file (use with --graph-format)",
    )

    parser.add_argument(
        "--graph-format", choices=["mermaid", "graphviz", "networkx"], default="mermaid",
        help="Graph export format: mermaid (default), graphviz (DOT), networkx (pickle/JSON)",
    )

    parser.add_argument(
        "--graph-stats", action="store_true",
        help="Show knowledge graph statistics (entities, relationships, types)",
    )

    parser.add_argument(
        "--deduplicate-kg", action="store_true",
        help="Deduplicate entities in knowledge graph (merge similar entity names)",
    )

    parser.add_argument(
        "--dedup-threshold", type=float, default=0.85,
        help="Similarity threshold for entity deduplication (default: 0.85, range: 0-1)",
    )

    parser.add_argument(
        "--dedup-report", action="store_true",
        help="Show deduplication report without applying changes",
    )

    parser.add_argument(
        "--graph-augmented", action="store_true",
        help="Enable graph-augmented retrieval (traverse KG edges for context expansion)",
    )

    parser.add_argument(
        "--graph-hops", type=int, default=2,
        help="Maximum graph traversal hops for graph-augmented retrieval (default: 2)",
    )

    # ── Collection arguments ──────────────────────────────────────────
    parser.add_argument(
        "--collection", "-c",
        metavar="NAME", action="append", default=[],
        help="Named collection to use (repeatable for blending: -c A -c B)",
    )

    parser.add_argument(
        "--add",
        metavar="PATH",
        help="Add documents at PATH to the specified collection(s)",
    )

    parser.add_argument(
        "--collection-list", action="store_true",
        help="List all available collections",
    )

    parser.add_argument(
        "--collection-info", action="store_true",
        help="Show detailed info for the specified collection",
    )

    parser.add_argument(
        "--collection-delete", action="store_true",
        help="Delete the specified collection and all its data",
    )

    parser.add_argument(
        "--collection-export",
        metavar="PATH",
        help="Export collection to a portable tar.gz archive at PATH (requires -c)",
    )

    parser.add_argument(
        "--collection-import",
        metavar="PATH",
        help="Import collection from a tar.gz archive at PATH",
    )

    parser.add_argument(
        "--import-rename",
        metavar="NAME",
        help="Rename imported collection to NAME (use with --collection-import)",
    )

    parser.add_argument(
        "--collection-merge",
        metavar="SRC",
        help="Merge SRC collection into target collection (requires -c for target)",
    )

    parser.add_argument(
        "--collection-description",
        metavar="TEXT", default="",
        help="Description for a new collection (used with --add)",
    )

    parser.add_argument(
        "--collection-tag",
        metavar="TAG",
        action="append", default=[],
        help="Add tag(s) to collection (use with -c) — repeatable",
    )

    parser.add_argument(
        "--collection-untag",
        metavar="TAG",
        action="append", default=[],
        help="Remove tag(s) from collection (use with -c) — repeatable",
    )

    parser.add_argument(
        "--collection-search",
        metavar="QUERY",
        help="Search collections by name/description (omit -c for all collections)",
    )

    parser.add_argument(
        "--collection-search-tags",
        metavar="TAGS",
        help="Filter collections by tags (comma-separated, use with --collection-search)",
    )

    parser.add_argument(
        "--collection-stats", action="store_true",
        help="Show detailed statistics for specified collection(s)",
    )

    parser.add_argument(
        "--global-stats", action="store_true",
        help="Show global statistics across all collections",
    )

    # Backward-compatible aliases (hidden from main help)
    parser.add_argument("--samba-nova", metavar="PATH", help=argparse.SUPPRESS)
    parser.add_argument("--nebius", metavar="PATH", help=argparse.SUPPRESS)

    args = parser.parse_args()

    # Handle --list-providers / --show-hierarchy
    if args.list_providers:
        list_providers()
        return

    if args.show_hierarchy:
        show_hierarchy()
        return

    if args.lock_status:
        status = lock_status()
        print(f"vrlmrag v{__version__} — Local Model Lock Status\n")
        if not status["locked"]:
            print("Status:  FREE — no local model is currently loaded in RAM")
        else:
            print("Status:  HELD — a local model is loaded in RAM")
            print(f"  PID:          {status['holder_pid']}")
            print(f"  Alive:        {status['holder_alive']}")
            print(f"  Model:        {status['model_id']}")
            print(f"  Since:        {status['acquired_at']}")
            print(f"  Process:      {status['process_name']}")
            print(f"  Description:  {status['description']}")
        print(f"\nThis process:   {'holds lock' if status['this_process_holds'] else 'does not hold lock'}")
        print(f"Lock file:      {Path.home() / '.vrlmrag' / 'local_model.lock'}")
        return

    # Resolve provider and input from args (support old-style flags)
    provider = args.provider
    input_path = args.input

    if args.samba_nova:
        provider = "sambanova"
        input_path = args.samba_nova
    elif args.nebius and not provider:
        provider = "nebius"
        input_path = args.nebius

    # ── Apply profile / comprehensive mode ──────────────────────────────
    if args.comprehensive:
        args.profile = "comprehensive"
    
    if args.profile:
        from vl_rag_graph_rlm.profiles import apply_profile
        apply_profile(args, args.profile)
        if not args.quiet:
            print(f"[profile] Using '{args.profile}' configuration preset")
            if args.profile == "comprehensive":
                print("[profile] Features enabled: multi-query, graph-augmented, verbose")

    # ── Compute use_api from --local and --offline flags ────────────────
    # Default: API mode (use_api=True). --local or --offline opts into local models.
    # --offline forces local mode and disables all API provider calls.
    use_api = not (args.local or args.offline)

    # ── Handle --offline mode ───────────────────────────────────────────
    if args.offline:
        # In offline mode, we cannot process video/audio at all
        if input_path:
            _input_ext = Path(input_path).suffix.lower()
            if _input_ext in _MEDIA_EXTENSIONS:
                print(f"[ERROR] Video/audio processing requires API mode.")
                print(f"        File: {input_path}")
                print(f"        Offline mode only supports text documents (TXT, MD, PPTX, PDF).")
                print(f"\n        To process this file, either:")
                print(f"          1. Remove --offline and use API mode (default)")
                print(f"          2. Set OPENROUTER_API_KEY and ZENMUX_API_KEY in .env")
                sys.exit(1)
        print("[offline] Running in offline mode — using local Qwen3-VL embeddings.")
        print("[offline] API providers disabled. Text documents only.")

    # Block local models for media files — always force API for video/audio
    if input_path and not use_api and not args.offline:
        _input_ext = Path(input_path).suffix.lower()
        if _input_ext in _MEDIA_EXTENSIONS:
            print(f"[safety] Local models are blocked for media files ({_input_ext}).")
            print(f"         Forcing API mode to prevent system overload.")
            use_api = True

    # ── Handle --reindex / --rebuild-kg flags ───────────────────────────
    if args.reindex or args.rebuild_kg:
        from pathlib import Path as _Path
        _store_path = _Path(input_path)
        if _store_path.is_file():
            _store_dir = _store_path.parent / ".vrlmrag_store"
        else:
            _store_dir = _store_path / ".vrlmrag_store"
        
        if args.reindex:
            # Clear embeddings to force re-processing
            _emb_file = _store_dir / "embeddings.json"
            _text_emb_file = _store_dir / "embeddings_text.json"
            _manifest_file = _store_dir / "manifest.json"
            if _emb_file.exists():
                _emb_file.unlink()
                print(f"[reindex] Cleared existing embeddings: {_emb_file}")
            if _text_emb_file.exists():
                _text_emb_file.unlink()
                print(f"[reindex] Cleared existing text embeddings: {_text_emb_file}")
            if _manifest_file.exists():
                _manifest_file.unlink()
                print(f"[reindex] Cleared manifest to force re-processing")
            print("[reindex] Will re-embed all documents with current model...")
        
        if args.rebuild_kg:
            # Clear knowledge graph to force re-extraction
            _kg_file = _store_dir / "knowledge_graph.md"
            if _kg_file.exists():
                _kg_file.unlink()
                print(f"[rebuild-kg] Cleared existing knowledge graph: {_kg_file}")
            print("[rebuild-kg] Will regenerate KG with current RLM model...")
        print()

    # ── Handle collection reindex / rebuild-kg ──────────────────────────
    if args.collection and (args.reindex or args.rebuild_kg):
        for cname in args.collection:
            if not collection_exists(cname):
                print(f"Error: Collection '{cname}' does not exist")
                print("Run 'vrlmrag --collection-list' to see available collections")
                sys.exit(1)
            
            meta = load_collection_meta(cname)
            slug = meta["name"]
            coll_dir = _collection_dir(slug)
            
            if args.reindex:
                # Clear collection embeddings
                _emb_file = coll_dir / "embeddings.json"
                _text_emb_file = coll_dir / "embeddings_text.json"
                if _emb_file.exists():
                    _emb_file.unlink()
                    print(f"[reindex:{slug}] Cleared embeddings")
                if _text_emb_file.exists():
                    _text_emb_file.unlink()
                    print(f"[reindex:{slug}] Cleared text embeddings")
            
            if args.rebuild_kg:
                # Clear collection KG
                _kg_file = coll_dir / "knowledge_graph.md"
                if _kg_file.exists():
                    _kg_file.unlink()
                    print(f"[rebuild-kg:{slug}] Cleared knowledge graph")
        
        print("\n[collection] Store cleared. Re-adding documents to rebuild...")
        # Now re-add all sources from the collection metadata
        for cname in args.collection:
            meta = load_collection_meta(cname)
            sources = meta.get("sources", [])
            if sources:
                print(f"\n[collection:{cname}] Re-adding {len(sources)} source(s)...")
                for src in sources:
                    src_path = src["path"]
                    if Path(src_path).exists():
                        # Run collection add with the same path
                        run_collection_add(
                            collection_names=[cname],
                            input_path=src_path,
                            provider=provider,
                            model=args.model,
                            max_depth=args.max_depth,
                            max_iterations=args.max_iterations,
                            description=meta.get("description", ""),
                            use_api=use_api,
                            text_only=args.text_only,
                        )
                    else:
                        print(f"  Warning: Source not found: {src_path}")
            else:
                print(f"[collection:{cname}] No sources recorded — nothing to reindex")
        return

    # ── Collection dispatch ────────────────────────────────────────────
    if args.quality_check:
        if not args.collection:
            print("Error: --quality-check requires -c <collection_name>")
            sys.exit(1)
        from vl_rag_graph_rlm.rlm_core import VLRAGGraphRLM
        for cname in args.collection:
            if not collection_exists(cname):
                print(f"Error: Collection '{cname}' does not exist")
                continue
            _run_quality_check(cname, provider, args.model, max_depth=args.max_depth, max_iterations=args.max_iterations)
        return

    # ── Handle graph export and stats ──────────────────────────────────
    if args.export_graph or args.graph_stats or args.deduplicate_kg or args.dedup_report:
        from vl_rag_graph_rlm.kg_visualization import (
            save_graph_visualization, get_graph_stats, 
            export_to_mermaid, export_to_graphviz
        )
        from vl_rag_graph_rlm.kg_deduplication import (
            deduplicate_knowledge_graph, get_deduplication_report
        )
        
        # Get knowledge graph from collection or input path
        kg_text = ""
        kg_source = ""
        if args.collection:
            for cname in args.collection:
                if collection_exists(cname):
                    from vl_rag_graph_rlm.collections import collection_load_kg
                    slug = load_collection_meta(cname)["name"]
                    kg = collection_load_kg(slug)
                    if kg:
                        kg_text += kg + "\n\n"
                        kg_source = f"collection:{cname}"
        elif input_path:
            store_path = Path(input_path)
            if store_path.is_file():
                store_path = store_path.parent
            kg_file = store_path / ".vrlmrag_store" / "knowledge_graph.md"
            if kg_file.exists():
                kg_text = kg_file.read_text(encoding="utf-8")
                kg_source = str(kg_file)
        
        if not kg_text.strip():
            print("Error: No knowledge graph found")
            print("Process documents first or specify a collection with -c")
            sys.exit(1)
        
        # Handle deduplication report
        if args.dedup_report:
            report = get_deduplication_report(kg_text, args.dedup_threshold)
            print(f"\nDeduplication Report (threshold: {args.dedup_threshold}):")
            print(f"  Total entities: {report['total_entities']}")
            print(f"  Duplicates found: {report['duplicates_found']}")
            if report['duplicate_pairs']:
                print("\n  Duplicate pairs:")
                for pair in report['duplicate_pairs'][:10]:  # Show top 10
                    print(f"    - '{pair['entity1']}' ≈ '{pair['entity2']}' ({pair['similarity']:.2f})")
        
        # Handle deduplication
        if args.deduplicate_kg:
            deduped_text, merge_map = deduplicate_knowledge_graph(kg_text, args.dedup_threshold)
            if merge_map:
                print(f"\nDeduplicated {len(merge_map)} entities:")
                for old, new in list(merge_map.items())[:10]:
                    print(f"  - '{old}' → '{new}'")
                if len(merge_map) > 10:
                    print(f"  ... and {len(merge_map) - 10} more")
                
                # Save back
                if args.collection:
                    for cname in args.collection:
                        if collection_exists(cname):
                            from vl_rag_graph_rlm.collections import collection_save_kg
                            slug = load_collection_meta(cname)["name"]
                            collection_save_kg(slug, deduped_text)
                            print(f"\nSaved deduplicated KG to collection: {cname}")
                elif input_path and kg_source:
                    Path(kg_source).write_text(deduped_text, encoding="utf-8")
                    print(f"\nSaved deduplicated KG to: {kg_source}")
                
                kg_text = deduped_text
            else:
                print("\nNo duplicates found at current threshold.")
        
        if args.graph_stats:
            stats = get_graph_stats(kg_text)
            print("\nKnowledge Graph Statistics:")
            print(f"  Entities: {stats['entity_count']}")
            print(f"  Relationships: {stats['relationship_count']}")
            print(f"  Connected ratio: {stats['connected_ratio']:.2f}")
            if stats['entity_types']:
                print("  Entity types:")
                for etype, count in sorted(stats['entity_types'].items(), key=lambda x: -x[1]):
                    print(f"    - {etype}: {count}")
        
        if args.export_graph:
            save_graph_visualization(kg_text, args.export_graph, args.graph_format)
            print(f"\nExported graph to: {args.export_graph} ({args.graph_format})")
        
        return

    if args.check_model:
        if not args.collection:
            print("Error: --check-model requires -c <collection_name>")
            sys.exit(1)
        from vl_rag_graph_rlm.collections import check_model_compatibility
        for cname in args.collection:
            if not collection_exists(cname):
                print(f"Error: Collection '{cname}' does not exist")
                continue
            result = check_model_compatibility(cname, args.check_model)
            print(f"\n[collection:{cname}] Model compatibility check:")
            print(f"  Current model: {result['current_model'] or '(not set)'}")
            print(f"  Target model:  {result['target_model']}")
            print(f"  Compatible:    {'Yes' if result['compatible'] else 'No'}")
            if result['needs_reindex']:
                print(f"  Action needed: Reindex required — models differ")
                print(f"    Run: vrlmrag -c {cname} --reindex")
            if result['mixed_models']:
                print(f"  Warning: Collection has mixed models: {', '.join(result['source_models'])}")
            if result['history']:
                print(f"  Model changes: {len(result['history'])} recorded")
        return

    if args.collection_list:
        show_collection_list()
        return

    if args.collection_delete:
        if not args.collection:
            print("Error: --collection-delete requires -c <name>")
            sys.exit(1)
        for cname in args.collection:
            if delete_collection(cname):
                print(f"Deleted collection: {cname}")
            else:
                print(f"Collection not found: {cname}")
        return

    # ── Handle collection export ───────────────────────────────────────
    if args.collection_export:
        if not args.collection:
            print("Error: --collection-export requires -c <name>")
            sys.exit(1)
        from vl_rag_graph_rlm.collections import export_collection
        for cname in args.collection:
            if not collection_exists(cname):
                print(f"Error: Collection '{cname}' does not exist")
                continue
            try:
                archive_path = export_collection(cname, args.collection_export)
                print(f"Exported collection '{cname}' to: {archive_path}")
            except Exception as e:
                print(f"Error exporting collection '{cname}': {e}")
        return

    # ── Handle collection import ───────────────────────────────────────
    if args.collection_import:
        from vl_rag_graph_rlm.collections import import_collection
        try:
            meta = import_collection(args.collection_import, args.import_rename)
            print(f"Imported collection as: {meta['name']}")
            print(f"  Display name: {meta.get('display_name', meta['name'])}")
            print(f"  Documents: {meta.get('document_count', 0)}")
            print(f"  Chunks: {meta.get('chunk_count', 0)}")
        except Exception as e:
            print(f"Error importing collection: {e}")
        return

    # ── Handle collection merge ────────────────────────────────────────
    if args.collection_merge:
        if not args.collection:
            print("Error: --collection-merge requires -c <target_name>")
            sys.exit(1)
        from vl_rag_graph_rlm.collections import merge_collections
        target_name = args.collection[0]
        if not collection_exists(target_name):
            print(f"Error: Target collection '{target_name}' does not exist")
            sys.exit(1)
        if not collection_exists(args.collection_merge):
            print(f"Error: Source collection '{args.collection_merge}' does not exist")
            sys.exit(1)
        try:
            meta = merge_collections(args.collection_merge, target_name)
            print(f"Merged '{args.collection_merge}' into '{target_name}'")
            print(f"  Total documents: {meta.get('document_count', 0)}")
            print(f"  Total chunks: {meta.get('chunk_count', 0)}")
        except Exception as e:
            print(f"Error merging collections: {e}")
        return

    if args.collection_info:
        if not args.collection:
            print("Error: --collection-info requires -c <name>")
            sys.exit(1)
        for cname in args.collection:
            show_collection_info(cname)
        return

    # ── Handle collection tagging ────────────────────────────────────
    if args.collection_tag:
        if not args.collection:
            print("Error: --collection-tag requires -c <name>")
            sys.exit(1)
        from vl_rag_graph_rlm.collections import add_tags
        for cname in args.collection:
            if not collection_exists(cname):
                print(f"Error: Collection '{cname}' does not exist")
                continue
            add_tags(cname, args.collection_tag)
            print(f"Added tags to '{cname}': {', '.join(args.collection_tag)}")
        return

    if args.collection_untag:
        if not args.collection:
            print("Error: --collection-untag requires -c <name>")
            sys.exit(1)
        from vl_rag_graph_rlm.collections import remove_tags
        for cname in args.collection:
            if not collection_exists(cname):
                print(f"Error: Collection '{cname}' does not exist")
                continue
            remove_tags(cname, args.collection_untag)
            print(f"Removed tags from '{cname}': {', '.join(args.collection_untag)}")
        return

    # ── Handle collection search ──────────────────────────────────────
    if args.collection_search or args.collection_search_tags:
        from vl_rag_graph_rlm.collections import search_collections
        search_tags = args.collection_search_tags.split(",") if args.collection_search_tags else None
        results = search_collections(
            query=args.collection_search,
            tags=search_tags,
        )
        if results:
            print(f"\nFound {len(results)} collection(s):")
            for c in results:
                tags = c.get("tags", [])
                tag_str = f" [{', '.join(tags)}]" if tags else ""
                print(f"  - {c.get('display_name', c['name'])}{tag_str}")
                print(f"    Documents: {c.get('document_count', 0)} | Chunks: {c.get('chunk_count', 0)}")
        else:
            print("No collections found matching criteria.")
        return

    # ── Handle collection statistics ──────────────────────────────────
    if args.collection_stats:
        from vl_rag_graph_rlm.collections import get_collection_stats
        if not args.collection:
            print("Error: --collection-stats requires -c <name>")
            sys.exit(1)
        for cname in args.collection:
            if not collection_exists(cname):
                print(f"Error: Collection '{cname}' does not exist")
                continue
            stats = get_collection_stats(cname)
            print(f"\n[collection:{cname}] Statistics:")
            print(f"  Documents: {stats['document_count']}")
            print(f"  Chunks: {stats['chunk_count']}")
            print(f"  Sources: {stats['sources_count']}")
            if stats['tags']:
                print(f"  Tags: {', '.join(stats['tags'])}")
            print(f"  Embedding model: {stats['embedding_model']}")
            print(f"  Knowledge graph: {stats['knowledge_graph_size']:,} chars")
            if stats['age_days'] is not None:
                print(f"  Age: {stats['age_days']} days")
        return

    if args.global_stats:
        from vl_rag_graph_rlm.collections import get_global_stats
        stats = get_global_stats()
        print("\nGlobal Collection Statistics:")
        print(f"  Total collections: {stats['total_collections']}")
        print(f"  Total documents: {stats['total_documents']}")
        print(f"  Total chunks: {stats['total_chunks']}")
        print(f"  Average docs/collection: {stats['average_documents']:.1f}")
        if stats['tag_distribution']:
            print("  Tag distribution:")
            for tag, count in list(stats['tag_distribution'].items())[:10]:
                print(f"    - {tag}: {count}")
        return

    if args.collection and args.add:
        run_collection_add(
            collection_names=args.collection,
            input_path=args.add,
            provider=provider,
            model=args.model,
            max_depth=args.max_depth,
            max_iterations=args.max_iterations,
            description=args.collection_description,
            use_api=use_api,
            text_only=args.text_only,
        )
        return

    if args.collection and args.query:
        run_collection_query(
            collection_names=args.collection,
            query=args.query,
            provider=provider,
            model=args.model,
            max_depth=args.max_depth,
            max_iterations=args.max_iterations,
            output=args.output,
            use_api=use_api,
            text_only=args.text_only,
        )
        return

    if args.collection and args.interactive:
        # Interactive mode with collection(s) as the store
        # Use the first collection's store dir for the session
        if not collection_exists(args.collection[0]):
            create_collection(args.collection[0])
        slug = load_collection_meta(args.collection[0])["name"]
        coll_dir = str(_collection_dir(slug))
        run_interactive_session(
            provider=provider,
            input_path=input_path,
            model=args.model,
            max_depth=args.max_depth,
            max_iterations=args.max_iterations,
            store_dir=coll_dir,
            use_api=use_api,
            text_only=args.text_only,
        )
        return

    # ── Interactive mode (no collection) ───────────────────────────────
    if args.interactive:
        run_interactive_session(
            provider=provider,
            input_path=input_path,
            model=args.model,
            max_depth=args.max_depth,
            max_iterations=args.max_iterations,
            store_dir=args.store_dir,
            use_api=use_api,
            text_only=args.text_only,
        )
        return

    # ── Handle --model-compare ─────────────────────────────────────────
    if args.model_compare and input_path:
        print(f"\n[model-compare] Comparing embeddings between '{args.model_compare}' and current model")
        print(f"[model-compare] Input: {input_path}")
        if args.query:
            print(f"[model-compare] Query: {args.query}")
        print()
        
        # Run with old model
        print(f"[model-compare] Running with OLD model: {args.model_compare}")
        print("-" * 50)
        old_result = run_analysis(
            provider=provider,
            input_path=input_path,
            query=args.query,
            output=None,  # Don't save output for comparison runs
            model=args.model_compare,
            max_depth=args.max_depth,
            max_iterations=args.max_iterations,
            use_api=use_api,
            text_only=args.text_only,
            _quiet=True,  # Suppress normal output
        )
        
        # Run with current model
        print(f"\n[model-compare] Running with CURRENT model: {args.model or 'default'}")
        print("-" * 50)
        new_result = run_analysis(
            provider=provider,
            input_path=input_path,
            query=args.query,
            output=None,
            model=args.model,
            max_depth=args.max_depth,
            max_iterations=args.max_iterations,
            use_api=use_api,
            text_only=args.text_only,
            _quiet=True,
        )
        
        # Compare and display results
        print("\n" + "=" * 60)
        print("MODEL COMPARISON RESULTS")
        print("=" * 60)
        
        from difflib import unified_diff
        old_lines = str(old_result.get('response', old_result)).split('\n')
        new_lines = str(new_result.get('response', new_result)).split('\n')
        
        diff = list(unified_diff(
            old_lines, new_lines,
            fromfile=f"OLD: {args.model_compare}",
            tofile=f"NEW: {args.model or 'default'}",
            lineterm=''
        ))
        
        if diff:
            print("\nDifferences detected:")
            for line in diff[:100]:  # Limit output
                print(line)
            if len(diff) > 100:
                print(f"... ({len(diff) - 100} more lines)")
        else:
            print("\nNo differences detected — responses are identical.")
        
        # Summary stats
        print(f"\nOld model response length: {len(str(old_result.get('response', '')))} chars")
        print(f"New model response length: {len(str(new_result.get('response', '')))} chars")
        print(f"Difference: {len(str(new_result.get('response', ''))) - len(str(old_result.get('response', '')))} chars")
        
        return

    # ── Default: run_analysis ──────────────────────────────────────────
    if not input_path:
        parser.print_help()
        print(f"\nError: input PATH is required.")
        sys.exit(1)

    run_analysis(
        provider=provider,
        input_path=input_path,
        query=args.query,
        output=args.output,
        model=args.model,
        max_depth=args.max_depth,
        max_iterations=args.max_iterations,
        use_api=use_api,
        text_only=args.text_only,
        multi_query=args.multi_query,
        rrf_weights=[args.rrf_dense_weight, args.rrf_keyword_weight],
        use_graph_augmented=args.graph_augmented,
        graph_hops=args.graph_hops,
        output_format=args.format,
        _quiet=args.quiet,
        verbose=args.verbose,
    )


if __name__ == "__main__":
    main()
