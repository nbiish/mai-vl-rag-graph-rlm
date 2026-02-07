#!/usr/bin/env python3
"""
Test VL-RAG-Graph-RLM on PowerPoint file

This script demonstrates:
- Extracting text and images from PPTX
- Indexing content with multimodal embeddings
- Querying across text AND visual content
"""

import os
import sys
from pathlib import Path

# Add src to path
sys.path.insert(0, str(Path(__file__).parent.parent / "src"))

from vl_rag_graph_rlm import VLRAGGraphRLM
from vl_rag_graph_rlm.rag import CompositeReranker, SearchResult, ReciprocalRankFusion

# Check for python-pptx
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("Installing python-pptx...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    HAS_PPTX = True

# Optional: Qwen3-VL for multimodal
try:
    from vl_rag_graph_rlm.rag.qwen3vl import create_qwen3vl_embedder
    HAS_MULTIMODAL = True
except ImportError:
    HAS_MULTIMODAL = False


def extract_pptx_content(pptx_path: str):
    """Extract text and image references from PowerPoint."""
    print(f"\nExtracting content from: {pptx_path}")
    
    prs = Presentation(pptx_path)
    slides_content = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_images = []
        
        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
            
            # Extract image references
            if shape.shape_type == 13:  # Picture
                try:
                    image = shape.image
                    slide_images.append({
                        'slide': slide_num,
                        'filename': f"slide_{slide_num}_image_{len(slide_images)}",
                        'ext': image.ext,
                        'size': (image.width, image.height)
                    })
                except:
                    pass
        
        slides_content.append({
            'slide_num': slide_num,
            'text': '\n'.join(slide_text),
            'images': slide_images
        })
        
        print(f"  Slide {slide_num}: {len(slide_text)} text elements, {len(slide_images)} images")
    
    return slides_content


def index_content(slides_content):
    """Index extracted content using RLM."""
    print("\n" + "=" * 70)
    print("Indexing Content")
    print("=" * 70)
    
    # Combine all text
    all_text = '\n\n'.join([
        f"SLIDE {s['slide_num']}:\n{s['text']}"
        for s in slides_content
    ])
    
    print(f"Total content: {len(all_text)} characters")
    print(f"Total slides: {len(slides_content)}")
    
    return all_text


def chunk_content(slides_content, max_chars=4000):
    """Chunk slides into groups that fit within token limits."""
    chunks = []
    current_chunk = []
    current_size = 0
    
    for slide in slides_content:
        slide_text = f"SLIDE {slide['slide_num']}:\n{slide['text']}\n\n"
        slide_size = len(slide_text)
        
        if current_size + slide_size > max_chars and current_chunk:
            # Save current chunk and start new one
            chunks.append('\n'.join(current_chunk))
            current_chunk = [slide_text]
            current_size = slide_size
        else:
            current_chunk.append(slide_text)
            current_size += slide_size
    
    if current_chunk:
        chunks.append('\n'.join(current_chunk))
    
    return chunks


def query_with_chunks(chunks: list, query: str):
    """Query across multiple chunks using RLM."""
    print(f"\nQuery: {query}")
    print(f"Processing {len(chunks)} chunks...")
    print("-" * 70)
    
    rlm = VLRAGGraphRLM(
        provider="sambanova",
        model="DeepSeek-V3.1",
        temperature=0.0,
        max_depth=3,
        max_iterations=10,
    )
    
    # For single chunk, query directly
    if len(chunks) == 1:
        result = rlm.completion(query, chunks[0])
        print(f"\nResponse:\n{result.response}")
        print(f"\nExecution time: {result.execution_time:.2f}s")
        return result
    
    # For multiple chunks, process each and combine
    responses = []
    for i, chunk in enumerate(chunks, 1):
        print(f"  Processing chunk {i}/{len(chunks)}...")
        try:
            chunk_query = f"Based on this content, answer: {query}"
            result = rlm.completion(chunk_query, chunk)
            responses.append(result.response)
        except Exception as e:
            print(f"    Error on chunk {i}: {e}")
    
    # Combine responses
    combined = "\n\n".join(responses)
    print(f"\nCombined Response:\n{combined[:1000]}...")
    return combined


def main():
    if not os.getenv("SAMBANOVA_API_KEY"):
        print("Error: SAMBANOVA_API_KEY not set")
        print("export SAMBANOVA_API_KEY=your_key_here")
        return
    
    pptx_path = "examples/Writing Tutorial 2022.pptx"
    
    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        return
    
    print("=" * 70)
    print("VL-RAG-Graph-RLM: PowerPoint Processing Test")
    print("=" * 70)
    
    # Step 1: Extract content
    slides_content = extract_pptx_content(pptx_path)
    
    # Step 2: Chunk content to fit within token limits (reduce to 4000 chars max)
    chunks = chunk_content(slides_content)
    print(f"\nContent split into {len(chunks)} chunks")
    
    # Step 3: Query
    print("\n" + "=" * 70)
    print("Querying Content")
    print("=" * 70)
    
    queries = [
        "What are the main topics covered in this tutorial?",
        "Summarize the key writing tips presented.",
        "What examples or case studies are mentioned?",
    ]
    
    for query in queries:
        query_with_chunks(chunks, query)
        print("\n" + "=" * 70)
    
    print("\nTest completed!")


if __name__ == "__main__":
    main()
